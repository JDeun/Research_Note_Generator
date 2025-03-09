import os
import json
import re
import base64
import logging
import time
import random
from typing import List, Dict, Any, Optional, Tuple, Union
from pathlib import Path
from datetime import datetime
import dotenv

# 설정 및 프롬프트 임포트
from config import DOCUMENT_PROCESSOR_MODELS, TEMPERATURES, LLM_API_KEY
from ProcessorPrompt import DOCUMENT_PROCESSOR_PROMPT

# langchain 관련 임포트
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.schema import Document, HumanMessage, SystemMessage
from langchain_community.callbacks import get_openai_callback

# PDF 처리 라이브러리
try:
    import pdfplumber
    PDFPLUMBER_AVAILABLE = True
except ImportError:
    PDFPLUMBER_AVAILABLE = False

try:
    import fitz  # PyMuPDF
    PYMUPDF_AVAILABLE = True
except ImportError:
    PYMUPDF_AVAILABLE = False

try:
    from pdfminer.high_level import extract_text, extract_pages
    from pdfminer.layout import LTTextContainer, LTFigure, LTImage, LTTextBox
    PDFMINER_AVAILABLE = True
except ImportError:
    PDFMINER_AVAILABLE = False

try:
    import PyPDF2
    PYPDF2_AVAILABLE = True
except ImportError:
    PYPDF2_AVAILABLE = False

try:
    import pypdfium2
    PYPDFIUM2_AVAILABLE = True
except ImportError:
    PYPDFIUM2_AVAILABLE = False

# DOCX 처리 라이브러리
try:
    import docx
    PYTHON_DOCX_AVAILABLE = True
except ImportError:
    PYTHON_DOCX_AVAILABLE = False

# PPTX 처리 라이브러리
try:
    from docling.document_converter import DocumentConverter
    DOCLING_AVAILABLE = True
except ImportError:
    DOCLING_AVAILABLE = False

try:
    from pptx import Presentation
    PYTHON_PPTX_AVAILABLE = True
except ImportError:
    PYTHON_PPTX_AVAILABLE = False

# XLSX 처리 라이브러리
try:
    import openpyxl
    OPENPYXL_AVAILABLE = True
except ImportError:
    OPENPYXL_AVAILABLE = False

# HTML 처리 라이브러리
try:
    from bs4 import BeautifulSoup
    BS4_AVAILABLE = True
except ImportError:
    BS4_AVAILABLE = False

# 로깅 설정
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# 문서 길이 임계값 설정
SHORT_DOC_THRESHOLD = 3000    # 3000자 미만은 짧은 문서로 간주
MEDIUM_DOC_THRESHOLD = 10000  # 3000-10000자는 중간 문서로 간주
# 10000자 이상은 긴 문서로 간주

# 청크 크기 설정 (문서 길이에 따라 다르게 적용)
CHUNK_SIZES = {
    "short": 2000,   # 짧은 문서
    "medium": 1500,  # 중간 문서
    "long": 1000     # 긴 문서
}

# 청크 겹침 설정
CHUNK_OVERLAPS = {
    "short": 200,   # 짧은 문서
    "medium": 150,  # 중간 문서
    "long": 100     # 긴 문서
}

# 최대 API 재시도 횟수
MAX_RETRIES = 5

# API 호출 간격 (초)
BASE_DELAY = 1.0

# LLM 템퍼러처 설정
TEMPERATURE = TEMPERATURES["document"]

class DocumentAnalyzer:
    """
    고급 문서 분석기
    - 문서 특성에 기반한 최적 파서 자동 선택
    - 복합 추출 기능 (텍스트, 표, 이미지)
    - 문서 구조 보존
    - 문서 길이에 따른 적응형 처리
    """
    
    def __init__(self, selected_model: str = "claude", auto_optimize: bool = True):
        """
        DocumentAnalyzer 초기화
        
        Args:
            selected_model (str): 사용할 LLM 모델명 (openai/gemini/claude/groq)
            auto_optimize (bool): 자동 파서 최적화 사용 여부
        """
        self.selected_model = selected_model.lower()
        self.auto_optimize = auto_optimize
        self.llm = None
        
        # API 호출 관리를 위한 카운터 초기화
        self.api_call_count = 0
        self.last_call_time = 0
        
        # 캐시 저장소 초기화
        self.summary_cache = {}
        
        # 설치된 파서 라이브러리 확인
        self.available_parsers = {
            "pdfminer": PDFMINER_AVAILABLE,
            "pdfplumber": PDFPLUMBER_AVAILABLE,
            "pymupdf": PYMUPDF_AVAILABLE,
            "pypdf2": PYPDF2_AVAILABLE,
            "pypdfium2": PYPDFIUM2_AVAILABLE,
            "docx": PYTHON_DOCX_AVAILABLE,
            "pptx": PYTHON_PPTX_AVAILABLE,
            "docling": DOCLING_AVAILABLE,
            "openpyxl": OPENPYXL_AVAILABLE,
            "bs4": BS4_AVAILABLE
        }
        
        # LLM 모델 설정
        self._setup_llm()
    
    def _setup_llm(self):
        """LLM 모델 설정"""
        if self.selected_model not in DOCUMENT_PROCESSOR_MODELS:
            logger.error(f"지원하지 않는 모델: {self.selected_model}")
            self.llm = None
            return

        api_model, model_class = DOCUMENT_PROCESSOR_MODELS[self.selected_model]
        api_key = LLM_API_KEY.get(self.selected_model)

        if not api_key:
            logger.error(f"API 키 누락: {self.selected_model} API 키를 .env에 설정하세요.")
            self.llm = None
            return

        try:
            self.llm = model_class(api_key=api_key, model=api_model, temperature=TEMPERATURE)
            logger.info(f"{self.selected_model} 모델 초기화 성공")
        except Exception as e:
            logger.error(f"{self.selected_model} 모델 초기화 실패: {str(e)}")
            self.llm = None
    
    def process_documents(self, file_paths: List[str]) -> List[Dict[str, Any]]:
        """
        여러 문서 파일 처리
        
        Args:
            file_paths (List[str]): 문서 파일 경로 목록
            
        Returns:
            List[Dict[str, Any]]: 처리 결과 목록
        """
        results = []
        for fp in file_paths:
            try:
                result = self.process_single_document(fp)
                results.append(result)
            except Exception as e:
                logger.error(f"문서 처리 실패 ({fp}): {str(e)}")
                results.append({"file_path": fp, "error": str(e)})
        return results
    
    def process_single_document(self, file_path: str) -> Dict[str, Any]:
        """
        단일 문서 파일 처리
        
        Args:
            file_path (str): 문서 파일 경로
            
        Returns:
            Dict[str, Any]: 처리 결과
        """
        try:
            # 파일 메타데이터 추출
            metadata = self._extract_metadata(file_path)
            file_type = metadata.get('file_type', 'Unknown')
            
            # 최적 파서 선택
            optimal_parser = self._select_optimal_parser(file_path, file_type)
            logger.info(f"선택된 최적 파서: {optimal_parser}")
            
            # 선택된 파서로 문서 로드 및 구조 추출
            docs, doc_structure = self._load_document_with_parser(file_path, file_type, optimal_parser)
            if not docs:
                return {
                    "file_info": metadata,
                    "error": "문서에서 텍스트를 추출하지 못했습니다."
                }
            
            # 전체 텍스트 생성
            full_text = "\n".join(d.page_content for d in docs)
            doc_length = len(full_text)
            
            # 문서 길이에 따른 처리 전략 결정
            if doc_length < SHORT_DOC_THRESHOLD:
                doc_size_category = "short"
            elif doc_length < MEDIUM_DOC_THRESHOLD:
                doc_size_category = "medium"
            else:
                doc_size_category = "long"
                
            logger.info(f"문서 길이: {doc_length}자 ({doc_size_category} 카테고리)")
            
            # 표 추출
            tables = []
            if file_type == "PDF" and self.available_parsers["pdfplumber"]:
                tables = self._extract_table_data(file_path)
                logger.info(f"추출된 표: {len(tables)}개")
            
            # 이미지 추출
            images = []
            if file_type == "PDF" and self.available_parsers["pymupdf"]:
                images = self._extract_images_from_pdf(file_path)
                logger.info(f"추출된 이미지: {len(images)}개")
            
            # 문서 길이에 따른 청크 크기 및 겹침 설정
            chunk_size = CHUNK_SIZES[doc_size_category]
            chunk_overlap = CHUNK_OVERLAPS[doc_size_category]
            
            # 텍스트 분할
            text_splitter = RecursiveCharacterTextSplitter(
                chunk_size=chunk_size,
                chunk_overlap=chunk_overlap
            )
            texts = text_splitter.split_text(full_text)
            
            # 처리할 청크 선택 (문서 길이에 따라 다르게)
            chunks_to_process = self._select_chunks_by_doc_size(texts, doc_size_category)
            
            # 선택된 청크만 요약
            summaries = self._process_chunks(chunks_to_process)
            
            # 요약 통합
            final_summary = self._generate_final_summary(summaries, full_text)
                
            # JSON 분석 결과 생성
            info_json = self._generate_info_json(final_summary, doc_length < SHORT_DOC_THRESHOLD)
            
            # 페이지 구조화 (한 페이지당 미리보기 생성)
            structured_text = self._structure_text(docs)
            
            # 최종 결과 구성
            result = {
                "file_info": metadata,
                "doc_length": doc_length,
                "doc_structure": doc_structure,
                "final_text_or_summary": final_summary,
                "tables": tables,
                "images": images,
                "text_content": structured_text,
                "analysis_result": info_json
            }
            
            return result
            
        except Exception as e:
            logger.error(f"문서 처리 중 오류 발생: {str(e)}")
            # 최소한의 결과 반환
            metadata = self._extract_metadata(file_path)
            return {
                "file_info": metadata,
                "error": str(e)
            }
    
    def _select_optimal_parser(self, file_path: str, file_type: str) -> str:
        """
        문서 유형 및 특성에 따라 최적의 파서 선택
        
        Args:
            file_path (str): 파일 경로
            file_type (str): 파일 유형
            
        Returns:
            str: 최적 파서 이름
        """
        # 기본 파서 설정
        default_parsers = {
            "PDF": "pdfminer",
            "DOCX": "docx",
            "PPTX": "pptx",
            "XLSX": "openpyxl",
            "TXT": "text"
        }
        
        if not self.auto_optimize:
            parser = default_parsers.get(file_type, "text")
            return parser if self.available_parsers.get(parser.lower(), False) else "text"
        
        # 파일 크기 및 초기 특성 분석
        file_size = os.path.getsize(file_path)
        
        if file_type == "PDF":
            # PDF 파일 특성 분석
            is_scanned, has_tables, has_images = self._analyze_pdf_features(file_path)
            
            # 특성에 따른 최적 파서 선택
            if is_scanned:
                # 스캔된 문서는 OCR이 필요하지만 현재 구현에는 OCR이 없음
                # 텍스트 레이어가 있을 수 있으므로 pdfminer로 시도
                if self.available_parsers["pdfminer"]:
                    return "pdfminer"
            
            if has_tables:
                # 표가 있는 문서는 pdfplumber가 좋음
                if self.available_parsers["pdfplumber"]:
                    return "pdfplumber"
            
            if has_images:
                # 이미지가 있는 문서는 pymupdf가 좋음
                if self.available_parsers["pymupdf"]:
                    return "pymupdf"
            
            # 기본값: 가용한 파서 중 선택
            for parser in ["pdfminer", "pdfplumber", "pymupdf", "pypdfium2", "pypdf2"]:
                if self.available_parsers[parser]:
                    return parser
        
        elif file_type == "DOCX":
            # DOCX 파일 최적 파서
            if self.available_parsers["docx"]:
                return "docx"
        
        elif file_type == "PPTX":
            # PPTX 파일 최적 파서
            if self.available_parsers["docling"]:
                return "docling"
            elif self.available_parsers["pptx"]:
                return "pptx"
        
        elif file_type == "XLSX":
            # XLSX 파일 최적 파서
            if self.available_parsers["openpyxl"]:
                return "openpyxl"
        
        # 기본값: 텍스트 파서
        return default_parsers.get(file_type, "text")
    
    def _analyze_pdf_features(self, file_path: str) -> Tuple[bool, bool, bool]:
        """
        PDF 파일의 특성 분석: 스캔 여부, 표 포함 여부, 이미지 포함 여부
        
        Args:
            file_path (str): PDF 파일 경로
            
        Returns:
            Tuple[bool, bool, bool]: (is_scanned, has_tables, has_images)
        """
        is_scanned = False
        has_tables = False
        has_images = False
        
        # PyMuPDF로 PDF 특성 확인
        if self.available_parsers["pymupdf"]:
            try:
                doc = fitz.open(file_path)
                
                # 처음 몇 페이지만 확인
                max_pages = min(5, len(doc))
                
                for page_idx in range(max_pages):
                    page = doc[page_idx]
                    
                    # 이미지 확인
                    image_list = page.get_images(full=True)
                    if len(image_list) > 0:
                        has_images = True
                    
                    # 텍스트 확인 (텍스트가 적으면 스캔 문서일 가능성)
                    text = page.get_text()
                    if len(text.strip()) < 100 and len(image_list) > 0:
                        is_scanned = True
                    
                    # 표 확인 (휴리스틱 방법: 표 관련 패턴 검색)
                    # 실제로는 더 정교한 표 감지 알고리즘 필요
                    if "표 " in text or "Table " in text or "|" in text:
                        has_tables = True
                
                doc.close()
                
            except Exception as e:
                logger.warning(f"PyMuPDF로 PDF 특성 분석 실패: {str(e)}")
        
        # PyPDF2로 특성 확인 (PyMuPDF 실패 시)
        elif self.available_parsers["pypdf2"]:
            try:
                with open(file_path, 'rb') as file:
                    reader = PyPDF2.PdfReader(file)
                    
                    # 처음 몇 페이지만 확인
                    max_pages = min(5, len(reader.pages))
                    
                    for page_idx in range(max_pages):
                        page = reader.pages[page_idx]
                        text = page.extract_text()
                        
                        # 텍스트가 적으면 스캔 문서일 가능성
                        if len(text.strip()) < 100:
                            is_scanned = True
                        
                        # 표 확인 (휴리스틱 방법)
                        if "표 " in text or "Table " in text or "|" in text:
                            has_tables = True
                            
            except Exception as e:
                logger.warning(f"PyPDF2로 PDF 특성 분석 실패: {str(e)}")
        
        # PDFPlumber로 표 확인
        if self.available_parsers["pdfplumber"]:
            try:
                with pdfplumber.open(file_path) as pdf:
                    # 처음 몇 페이지만 확인
                    max_pages = min(5, len(pdf.pages))
                    
                    for i in range(max_pages):
                        tables = pdf.pages[i].extract_tables()
                        if tables and len(tables) > 0:
                            has_tables = True
                            break
                            
            except Exception as e:
                logger.warning(f"PDFPlumber로 표 확인 실패: {str(e)}")
        
        return is_scanned, has_tables, has_images
    
    def _load_document_with_parser(self, file_path: str, file_type: str, parser: str) -> Tuple[List[Document], Dict[str, Any]]:
        """
        선택된 파서로 문서 로드
        
        Args:
            file_path (str): 파일 경로
            file_type (str): 파일 유형
            parser (str): 파서 이름
            
        Returns:
            Tuple[List[Document], Dict[str, Any]]: (문서 리스트, 구조 정보)
        """
        docs = []
        structure_info = {"type": file_type, "parser": parser}
        
        try:
            logger.info(f"파서 '{parser}'로 '{file_type}' 파일 로드 시작: {file_path}")
            
            # PDF 파일 처리
            if file_type == "PDF":
                if parser == "pdfminer":
                    docs, structure_info = self._load_with_pdfminer(file_path)
                elif parser == "pdfplumber":
                    docs, structure_info = self._load_with_pdfplumber(file_path)
                elif parser == "pymupdf":
                    docs, structure_info = self._load_with_pymupdf(file_path)
                elif parser == "pypdf2":
                    docs, structure_info = self._load_with_pypdf2(file_path)
                elif parser == "pypdfium2":
                    docs, structure_info = self._load_with_pypdfium2(file_path)
                else:
                    # 기본 파서로 폴백
                    docs, structure_info = self._load_with_text(file_path)
            
            # DOCX 파일 처리
            elif file_type == "DOCX":
                if parser == "docx" and self.available_parsers["docx"]:
                    docs, structure_info = self._load_with_docx(file_path)
                else:
                    # 텍스트 추출 폴백
                    docs, structure_info = self._load_with_text(file_path)
            
            # PPTX 파일 처리
            elif file_type == "PPTX":
                if parser == "docling" and self.available_parsers["docling"]:
                    docs, structure_info = self._load_with_docling(file_path)
                elif parser == "pptx" and self.available_parsers["pptx"]:
                    docs, structure_info = self._load_with_pptx(file_path)
                else:
                    # 텍스트 추출 폴백
                    docs, structure_info = self._load_with_text(file_path)
            
            # XLSX 파일 처리
            elif file_type == "XLSX":
                if parser == "openpyxl" and self.available_parsers["openpyxl"]:
                    docs, structure_info = self._load_with_openpyxl(file_path)
                else:
                    # 텍스트 추출 폴백
                    docs, structure_info = self._load_with_text(file_path)
            
            # 기타 텍스트 파일 처리
            else:
                docs, structure_info = self._load_with_text(file_path)
            
            # 문서가 없으면 빈 문서 생성
            if not docs:
                docs = [Document(page_content="문서에서 텍스트를 추출할 수 없습니다.", metadata={"source": file_path})]
            
            logger.info(f"문서 로드 완료: {len(docs)}개 문서 조각")
            return docs, structure_info
            
        except Exception as e:
            logger.error(f"문서 로드 실패 ({parser}): {str(e)}")
            # 기본 문서 반환
            return [Document(page_content=f"문서 로드 중 오류 발생: {str(e)}", metadata={"source": file_path})], {"error": str(e)}
    
    def _load_with_pdfminer(self, file_path: str) -> Tuple[List[Document], Dict[str, Any]]:
        """PDFMiner를 사용한 PDF 로드"""
        try:
            structure_info = {"type": "PDF", "parser": "pdfminer"}
            docs = []
            
            # 페이지별 텍스트 추출
            for i, page_layout in enumerate(extract_pages(file_path)):
                page_text = ""
                
                # 레이아웃 요소 처리
                for element in page_layout:
                    if isinstance(element, LTTextContainer):
                        page_text += element.get_text()
                
                # 페이지 문서 생성
                if page_text.strip():
                    metadata = {"source": file_path, "page": i + 1}
                    docs.append(Document(page_content=page_text.strip(), metadata=metadata))
            
            # 구조 정보 업데이트
            structure_info["total_pages"] = len(docs)
            
            return docs, structure_info
            
        except Exception as e:
            logger.error(f"PDFMiner 로드 실패: {str(e)}")
            return [], {"error": str(e)}
    
    def _load_with_pdfplumber(self, file_path: str) -> Tuple[List[Document], Dict[str, Any]]:
        """PDFPlumber를 사용한 PDF 로드"""
        try:
            structure_info = {"type": "PDF", "parser": "pdfplumber", "pages": []}
            docs = []
            
            with pdfplumber.open(file_path) as pdf:
                structure_info["total_pages"] = len(pdf.pages)
                
                for i, page in enumerate(pdf.pages):
                    page_text = page.extract_text() or ""
                    
                    # 페이지 정보 저장
                    page_info = {
                        "page_num": i + 1,
                        "width": page.width,
                        "height": page.height,
                        "text_length": len(page_text)
                    }
                    structure_info["pages"].append(page_info)
                    
                    # 문서 객체 생성
                    metadata = {"source": file_path, "page": i + 1, "width": page.width, "height": page.height}
                    docs.append(Document(page_content=page_text, metadata=metadata))
            
            return docs, structure_info
            
        except Exception as e:
            logger.error(f"PDFPlumber 로드 실패: {str(e)}")
            return [], {"error": str(e)}
    
    def _load_with_pymupdf(self, file_path: str) -> Tuple[List[Document], Dict[str, Any]]:
        """PyMuPDF를 사용한 PDF 로드"""
        try:
            structure_info = {"type": "PDF", "parser": "pymupdf", "pages": []}
            docs = []
            
            doc = fitz.open(file_path)
            structure_info["total_pages"] = len(doc)
            
            for i in range(len(doc)):
                page = doc[i]
                page_text = page.get_text()
                
                # 페이지 정보
                page_info = {
                    "page_num": i + 1,
                    "width": page.rect.width,
                    "height": page.rect.height,
                    "text_length": len(page_text)
                }
                structure_info["pages"].append(page_info)
                
                # 문서 객체 생성
                metadata = {"source": file_path, "page": i + 1, "width": page.rect.width, "height": page.rect.height}
                docs.append(Document(page_content=page_text, metadata=metadata))
            
            doc.close()
            return docs, structure_info
            
        except Exception as e:
            logger.error(f"PyMuPDF 로드 실패: {str(e)}")
            return [], {"error": str(e)}
    
    def _load_with_pypdf2(self, file_path: str) -> Tuple[List[Document], Dict[str, Any]]:
        """PyPDF2를 사용한 PDF 로드"""
        try:
            structure_info = {"type": "PDF", "parser": "pypdf2", "pages": []}
            docs = []
            
            with open(file_path, 'rb') as f:
                reader = PyPDF2.PdfReader(f)
                structure_info["total_pages"] = len(reader.pages)
                
                for i, page in enumerate(reader.pages):
                    page_text = page.extract_text() or ""
                    
                    # 페이지 정보
                    page_info = {
                        "page_num": i + 1,
                        "text_length": len(page_text)
                    }
                    structure_info["pages"].append(page_info)
                    
                    # 문서 객체 생성
                    metadata = {"source": file_path, "page": i + 1}
                    docs.append(Document(page_content=page_text, metadata=metadata))
            
            return docs, structure_info
            
        except Exception as e:
            logger.error(f"PyPDF2 로드 실패: {str(e)}")
            return [], {"error": str(e)}
    
    def _load_with_pypdfium2(self, file_path: str) -> Tuple[List[Document], Dict[str, Any]]:
        """PyPDFium2를 사용한 PDF 로드"""
        try:
            structure_info = {"type": "PDF", "parser": "pypdfium2", "pages": []}
            docs = []
            
            pdf = pypdfium2.PdfDocument(file_path)
            structure_info["total_pages"] = len(pdf)
            
            for i in range(len(pdf)):
                page = pdf[i]
                textpage = page.get_textpage()
                page_text = textpage.get_text() or ""
                
                # 페이지 정보
                page_info = {
                    "page_num": i + 1,
                    "width": page.get_size()[0],
                    "height": page.get_size()[1],
                    "text_length": len(page_text)
                }
                structure_info["pages"].append(page_info)
                
                # 문서 객체 생성
                metadata = {"source": file_path, "page": i + 1}
                docs.append(Document(page_content=page_text, metadata=metadata))
                
                # 자원 해제
                textpage.close()
                page.close()
            
            pdf.close()
            return docs, structure_info
            
        except Exception as e:
            logger.error(f"PyPDFium2 로드 실패: {str(e)}")
            return [], {"error": str(e)}
    
    def _load_with_docx(self, file_path: str) -> Tuple[List[Document], Dict[str, Any]]:
        """python-docx를 사용한 DOCX 로드"""
        try:
            structure_info = {"type": "DOCX", "parser": "docx", "sections": []}
            docs = []
            
            # 문서 로드
            doc = docx.Document(file_path)
            
            # 기본 정보 추출
            structure_info["paragraphs"] = len(doc.paragraphs)
            structure_info["tables"] = len(doc.tables)
            
            # 제목 추출
            headings = []
            for i, para in enumerate(doc.paragraphs):
                if para.style.name.startswith('Heading'):
                    level = int(para.style.name.replace('Heading', '')) if para.style.name != 'Heading' else 1
                    headings.append({
                        "level": level,
                        "text": para.text,
                        "position": i
                    })
            structure_info["headings"] = headings
            
            # 섹션 단위로 분리 (제목 기준)
            sections = []
            current_section = {"title": "", "content": ""}
            
            for i, para in enumerate(doc.paragraphs):
                if para.style.name.startswith('Heading'):
                    # 기존 섹션 저장
                    if current_section["content"].strip():
                        sections.append(current_section)
                    
                    # 새 섹션 시작
                    current_section = {"title": para.text, "content": ""}
                else:
                    if para.text.strip():
                        current_section["content"] += para.text + "\n"
            
            # 마지막 섹션 저장
            if current_section["content"].strip():
                sections.append(current_section)
            
            structure_info["sections"] = sections
            
            # 섹션별 문서 생성
            for i, section in enumerate(sections):
                metadata = {
                    "source": file_path,
                    "section": i + 1,
                    "title": section["title"]
                }
                docs.append(Document(page_content=section["content"], metadata=metadata))
            
            # 섹션이 없으면 전체 내용을 하나의 문서로
            if not docs:
                full_text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
                metadata = {"source": file_path}
                docs.append(Document(page_content=full_text, metadata=metadata))
            
            return docs, structure_info
            
        except Exception as e:
            logger.error(f"DOCX 로드 실패: {str(e)}")
            return [], {"error": str(e)}
    
    def _load_with_docling(self, file_path: str) -> Tuple[List[Document], Dict[str, Any]]:
        """Docling을 사용한 PPTX 로드"""
        try:
            structure_info = {"type": "PPTX", "parser": "docling"}
            docs = []
            
            # Docling으로 변환
            converter = DocumentConverter()
            result = converter.convert(file_path)
            md_text = result.document.export_to_markdown()
            
            if not md_text.strip():
                logger.warning("Docling 변환 결과가 비어있습니다.")
                return [], {"error": "변환 결과 없음"}
            
            # 슬라이드 구분자 패턴
            slide_pattern = r'#{2,} Slide \d+'
            
            # 슬라이드별로 분리
            slide_texts = re.split(slide_pattern, md_text)
            slide_titles = re.findall(slide_pattern, md_text)
            
            # 첫 번째 빈 요소 제거
            if slide_texts and not slide_texts[0].strip():
                slide_texts = slide_texts[1:]
            
            # 구조 정보 업데이트
            structure_info["total_slides"] = len(slide_texts)
            structure_info["slides"] = []
            
            # 슬라이드별 문서 생성
            for i, (slide_text, slide_title) in enumerate(zip(slide_texts, slide_titles + [""] * (len(slide_texts) - len(slide_titles)))):
                if not slide_text.strip():
                    continue
                
                slide_info = {
                    "slide_num": i + 1,
                    "title": slide_title.strip() if slide_title else f"Slide {i+1}"
                }
                structure_info["slides"].append(slide_info)
                
                metadata = {
                    "source": file_path,
                    "slide": i + 1,
                    "title": slide_title.strip() if slide_title else f"Slide {i+1}"
                }
                docs.append(Document(page_content=slide_text.strip(), metadata=metadata))
            
            return docs, structure_info
            
        except Exception as e:
            logger.error(f"Docling 로드 실패: {str(e)}")
            return [], {"error": str(e)}
    
    def _load_with_pptx(self, file_path: str) -> Tuple[List[Document], Dict[str, Any]]:
        """python-pptx를 사용한 PPTX 로드"""
        try:
            structure_info = {"type": "PPTX", "parser": "pptx", "slides": []}
            docs = []
            
            # PPTX 로드
            prs = Presentation(file_path)
            
            # 구조 정보
            structure_info["total_slides"] = len(prs.slides)
            
            # 슬라이드별 처리
            for i, slide in enumerate(prs.slides):
                slide_text = []
                slide_info = {
                    "slide_num": i + 1,
                    "shapes": len(slide.shapes)
                }
                
                # 슬라이드 텍스트 추출
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text)
                
                # 슬라이드 노트 확인
                notes_text = ""
                if slide.has_notes_slide and slide.notes_slide.notes_text_frame.text.strip():
                    notes_text = slide.notes_slide.notes_text_frame.text
                    slide_info["has_notes"] = True
                
                # 구조 정보 업데이트
                combined_text = "\n".join(slide_text)
                if notes_text:
                    combined_text += "\n\n[Notes]\n" + notes_text
                
                slide_info["text_length"] = len(combined_text)
                structure_info["slides"].append(slide_info)
                
                # 문서 객체 생성
                metadata = {
                    "source": file_path,
                    "slide": i + 1,
                    "has_notes": bool(notes_text)
                }
                docs.append(Document(page_content=combined_text, metadata=metadata))
            
            return docs, structure_info
            
        except Exception as e:
            logger.error(f"PPTX 로드 실패: {str(e)}")
            return [], {"error": str(e)}
    
    def _load_with_openpyxl(self, file_path: str) -> Tuple[List[Document], Dict[str, Any]]:
        """openpyxl을 사용한 XLSX 로드"""
        try:
            structure_info = {"type": "XLSX", "parser": "openpyxl", "sheets": []}
            docs = []
            
            # XLSX 로드
            wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
            
            # 구조 정보
            structure_info["total_sheets"] = len(wb.sheetnames)
            structure_info["sheet_names"] = wb.sheetnames
            
            # 시트별 처리
            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                
                # 시트 정보
                sheet_info = {
                    "name": sheet_name,
                    "max_row": sheet.max_row,
                    "max_column": sheet.max_column
                }
                structure_info["sheets"].append(sheet_info)
                
                # 시트 내용 추출
                sheet_content = []
                for row in sheet.iter_rows(values_only=True):
                    if any(cell is not None and str(cell).strip() for cell in row):
                        sheet_content.append("\t".join(str(cell) if cell is not None else "" for cell in row))
                
                # 문서 객체 생성
                sheet_text = "\n".join(sheet_content)
                metadata = {
                    "source": file_path,
                    "sheet": sheet_name,
                    "rows": sheet.max_row,
                    "columns": sheet.max_column
                }
                docs.append(Document(page_content=sheet_text, metadata=metadata))
            
            # 워크북 닫기
            wb.close()
            return docs, structure_info
            
        except Exception as e:
            logger.error(f"XLSX 로드 실패: {str(e)}")
            return [], {"error": str(e)}
    
    def _load_with_text(self, file_path: str) -> Tuple[List[Document], Dict[str, Any]]:
        """일반 텍스트 파일 로드"""
        try:
            structure_info = {"type": "TEXT", "parser": "text"}
            
            # 텍스트 로드 시도
            try:
                with open(file_path, 'r', encoding='utf-8') as f:
                    text = f.read()
            except UnicodeDecodeError:
                # UTF-8 실패 시 다른 인코딩 시도
                encodings = ['cp949', 'euc-kr', 'latin1']
                for encoding in encodings:
                    try:
                        with open(file_path, 'r', encoding=encoding) as f:
                            text = f.read()
                        break
                    except UnicodeDecodeError:
                        continue
                else:
                    # 모든 인코딩 실패 시
                    raise ValueError("지원되지 않는 파일 인코딩")
            
            # 구조 정보 업데이트
            lines = text.split('\n')
            paragraphs = [p for p in text.split('\n\n') if p.strip()]
            
            structure_info["lines"] = len(lines)
            structure_info["paragraphs"] = len(paragraphs)
            
            # 문서 객체 생성
            metadata = {"source": file_path}
            docs = [Document(page_content=text, metadata=metadata)]
            
            return docs, structure_info
            
        except Exception as e:
            logger.error(f"텍스트 로드 실패: {str(e)}")
            return [], {"error": str(e)}
        
    def _extract_page_as_image(self, pdf_path: str) -> List[Dict]:
        """PDF 페이지를 이미지로 렌더링하여 추출"""
        if not PYMUPDF_AVAILABLE:
            logger.warning("PyMuPDF 라이브러리가 설치되지 않았습니다.")
            return []
            
        images = []
        doc = fitz.open(pdf_path)
        
        for page_idx in range(len(doc)):
            page = doc[page_idx]
            # 고해상도로 페이지 렌더링
            pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
            img_data = pix.tobytes("png")
            
            images.append({
                "page": page_idx + 1,
                "type": "png",
                "caption": f"페이지 {page_idx+1} 이미지",
                "data": base64.b64encode(img_data).decode('utf-8')
            })
        
        doc.close()
        return images
    
    def _detect_graph_regions(self, page) -> List[Dict]:
        """PDF 페이지에서 그래프 영역을 감지"""
        graph_regions = []
        
        # 그래프 관련 키워드 패턴
        graph_patterns = ["그래프", "차트", "plot", "figure", "graph", "chart", "diagram"]
        
        # 페이지 텍스트 분석
        text_blocks = page.get_text("blocks")
        
        for block in text_blocks:
            text = block[4]
            rect = fitz.Rect(block[:4])
            
            # 텍스트에 그래프 관련 키워드가 있는지 확인
            if any(pattern in text.lower() for pattern in graph_patterns):
                # 키워드 아래 영역을 그래프 영역으로 가정
                graph_rect = fitz.Rect(rect.x0, rect.y1, rect.x1, rect.y1 + 200)
                
                # 해당 영역을 이미지로 추출
                pix = page.get_pixmap(clip=graph_rect, matrix=fitz.Matrix(2, 2))
                
                graph_regions.append({
                    "rect": graph_rect,
                    "caption": text.strip(),
                    "image": pix
                })
        
        return graph_regions
    
    def _classify_image_type(self, image_data, caption: str) -> str:
        """이미지 유형 분류 (그래프, 차트, 일반 이미지 등)"""
        # 캡션 기반 분류
        graph_keywords = ["그래프", "chart", "plot", "graph", "도표", "figure"]
        table_keywords = ["표", "table"]
        
        if any(keyword in caption.lower() for keyword in graph_keywords):
            return "graph"
        elif any(keyword in caption.lower() for keyword in table_keywords):
            return "table"
        else:
            return "image"
    
    def _select_chunks_by_doc_size(self, texts: List[str], doc_size_category: str) -> List[str]:
        """
        문서 크기별 처리할 청크 선택
        
        Args:
            texts (List[str]): 전체 청크 목록
            doc_size_category (str): 문서 크기 카테고리
            
        Returns:
            List[str]: 선택된 청크 목록
        """
        if not texts:
            return []
            
        if doc_size_category == "short" or len(texts) <= 3:
            # 짧은 문서는 모든 청크 처리
            return texts
        elif doc_size_category == "medium":
            # 중간 문서는 주요 부분만 처리
            if len(texts) <= 5:
                return texts
            else:
                # 첫 부분, 중간 부분, 마지막 부분 선택
                return [
                    texts[0],                   # 첫 번째 청크
                    texts[len(texts) // 4],     # 1/4 지점
                    texts[len(texts) // 2],     # 중간 지점
                    texts[3 * len(texts) // 4], # 3/4 지점
                    texts[-1]                   # 마지막 청크
                ]
        else:
            # 긴 문서는 더 제한적으로 처리
            if len(texts) <= 3:
                return texts
            else:
                # 첫 부분, 중간 부분, 마지막 부분만 처리
                return [
                    texts[0],                   # 첫 번째 청크
                    texts[len(texts) // 2],     # 중간 지점
                    texts[-1]                   # 마지막 청크
                ]
    
    def _process_chunks(self, chunks: List[str]) -> List[str]:
        """
        선택된 청크 처리 (요약 생성)
        
        Args:
            chunks (List[str]): 처리할 청크 목록
            
        Returns:
            List[str]: 청크별 요약 목록
        """
        if not chunks:
            logger.debug("처리할 청크가 없습니다.")
            return []
        
        summaries = []
        for i, chunk in enumerate(chunks):
            try:
                logger.debug(f"청크 {i+1}/{len(chunks)} 처리 중...")
                chunk_hash = hash(chunk)
                
                # 캐시에서 요약 확인
                if chunk_hash in self.summary_cache:
                    summary = self.summary_cache[chunk_hash]
                    logger.debug("캐시에서 요약 로드됨")
                else:
                    # 새로 요약 생성
                    try:
                        summary = self._generate_chunk_summary_with_retry(chunk)
                        if summary:
                            # 요약 캐싱
                            self.summary_cache[chunk_hash] = summary
                    except Exception as e:
                        import traceback
                        logger.error(f"청크 요약 생성 실패: {str(e)}\n{traceback.format_exc()}")
                        # 실패 시 청크 앞부분 사용
                        summary = chunk[:200] + "... (요약 생성 실패)"
                
                summaries.append(summary)
                
            except Exception as e:
                import traceback
                logger.error(f"청크 {i+1} 처리 중 오류 발생: {str(e)}\n{traceback.format_exc()}")
                # 오류 발생 시 빈 요약 추가
                summaries.append("요약 생성 중 오류 발생")
        
        return summaries
    
    def _manage_api_rate(self):
        """API 호출 속도 관리"""
        current_time = time.time()
        elapsed = current_time - self.last_call_time
        
        # 최소 0.5초 대기
        if elapsed < 0.5:
            time.sleep(0.5 - elapsed)
        
        # 호출 추적
        self.api_call_count += 1
        self.last_call_time = time.time()
    
    def _retry_with_backoff(self, func, args=None, kwargs=None, max_retries=MAX_RETRIES, initial_delay=BASE_DELAY):
        """지수 백오프를 적용한 재시도 로직"""
        args = args or []
        kwargs = kwargs or {}
        
        for attempt in range(max_retries):
            try:
                # API 호출 관리
                self._manage_api_rate()
                return func(*args, **kwargs)
            except Exception as e:
                # 할당량 초과 오류 확인
                is_rate_limit = "429" in str(e) or "quota" in str(e).lower() or "rate limit" in str(e).lower()
                
                if attempt < max_retries - 1:
                    # 백오프 시간 계산
                    wait_time = (2 ** attempt) * initial_delay + random.uniform(0, initial_delay)
                    
                    # 할당량 초과 시 더 오래 대기
                    if is_rate_limit:
                        wait_time *= 2
                        logger.warning(f"API 할당량 초과. {wait_time:.1f}초 대기 후 재시도 ({attempt+1}/{max_retries})")
                    else:
                        logger.warning(f"오류 발생. {wait_time:.1f}초 대기 후 재시도 ({attempt+1}/{max_retries})")
                    
                    time.sleep(wait_time)
                    continue
                else:
                    # 최대 재시도 초과
                    logger.error(f"최대 재시도 횟수 초과: {str(e)}")
                    raise
    
    def _generate_chunk_summary_with_retry(self, text: str) -> str:
        """재시도 로직이 적용된 청크 요약 생성"""
        return self._retry_with_backoff(self._generate_chunk_summary, args=[text])
    
    def _generate_chunk_summary(self, text: str) -> str:
        """청크 요약 생성"""
        if not self.llm:
            logger.warning("LLM이 초기화되지 않았습니다. 요약을 건너뜁니다.")
            return self._fallback_summary(text)
        
        try:
            # 요약 생성 프롬프트
            prompt = SystemMessage(content=DOCUMENT_PROCESSOR_PROMPT["system_summary"])
            user_prompt = HumanMessage(content=DOCUMENT_PROCESSOR_PROMPT["user_summary"].format(
                full_text=text
            ))
            
            # LLM 호출
            response = self.llm.invoke([prompt, user_prompt])
            return response.content
            
        except Exception as e:
            logger.error(f"청크 요약 생성 실패: {str(e)}")
            return self._fallback_summary(text)
    
    def _fallback_summary(self, text: str) -> str:
        """LLM 요약 실패 시 폴백 요약 생성"""
        # 간단한 규칙 기반 요약
        sentences = re.split(r'(?<=[.!?])\s+', text)
        
        if len(sentences) <= 5:
            return text
        
        # 첫 3문장과 마지막 2문장 추출
        summary_sentences = sentences[:3] + sentences[-2:]
        summary = ' '.join(summary_sentences)
        
        return summary + "\n(자동 추출 요약)"
    
    def _generate_final_summary(self, summaries: List[str], full_text: str) -> str:
        """최종 요약 생성"""
        # 요약이 없으면 전체 텍스트의 앞부분 사용
        if not summaries:
            logger.debug("요약이 없어 전체 텍스트 앞부분을 사용합니다.")
            return full_text[:min(SHORT_DOC_THRESHOLD, len(full_text))]
            
        # 단일 요약이면 그대로 반환
        if len(summaries) == 1:
            logger.debug("단일 요약을 그대로 반환합니다.")
            return summaries[0]
            
        # 요약 결합
        try:
            combined = "\n\n".join(summaries)
            
            # LLM이 없거나 결합된 요약이 짧으면 그대로 반환
            if not self.llm or len(combined) < 1000:
                logger.debug("LLM이 없거나 요약이 짧아 결합된 요약을 그대로 반환합니다.")
                return combined
            
            # 최종 요약 생성
            logger.debug("LLM을 사용하여 최종 요약 생성 중...")
            prompt = SystemMessage(content=DOCUMENT_PROCESSOR_PROMPT["system_summary"])
            user_prompt = HumanMessage(content=DOCUMENT_PROCESSOR_PROMPT["user_summary"].format(
                full_text=combined
            ))
            
            response = self.llm.invoke([prompt, user_prompt])
            logger.debug("최종 요약 생성 완료")
            return response.content
            
        except Exception as e:
            import traceback
            logger.error(f"요약 통합 실패: {str(e)}\n{traceback.format_exc()}")
            # 결합된 요약 반환
            if isinstance(summaries, list) and summaries:
                return "\n\n".join(summaries)  # 실패 시 요약 연결
            else:
                logger.warning("summaries가 리스트가 아니거나 비어 있습니다.")
                return str(summaries) if summaries else "요약을 생성할 수 없습니다."
    
    def _generate_info_json(self, summary: str, is_short: bool) -> Dict[str, str]:
        """문서 정보 JSON 생성"""
        if not self.llm:
            logger.debug("LLM이 초기화되지 않아 기본 정보 추출을 사용합니다.")
            return self._fallback_info_extraction(summary)
        
        try:
            # LLM 응답에 코드 블록 마크다운이 포함된 경우 이를 제거하는 전처리
            def cleanup_response(text: str) -> str:
                pattern = re.compile(r"```(?:json)?\s*([\s\S]*?)\s*```", re.MULTILINE)
                m = pattern.search(text)
                if m:
                    return m.group(1).strip()
                return text.strip()
            
            # 프롬프트 생성
            logger.debug(f"{'전체 문서' if is_short else '요약된 문서'}에 대한 정보 추출 중...")
            if is_short:
                system_prompt = SystemMessage(content=DOCUMENT_PROCESSOR_PROMPT["system_full"])
                user_prompt = HumanMessage(content=DOCUMENT_PROCESSOR_PROMPT["user_full"].format(final_text=summary))
            else:
                system_prompt = SystemMessage(content=DOCUMENT_PROCESSOR_PROMPT["system"])
                user_prompt = HumanMessage(content=DOCUMENT_PROCESSOR_PROMPT["user"].format(final_summary=summary))
            
            # LLM 호출
            resp = self.llm.invoke([system_prompt, user_prompt])
            content = cleanup_response(resp.content)
            logger.debug("LLM 응답 받음, JSON 변환 시도 중...")
            
            try:
                # 응답 파싱
                data = json.loads(content)
                
                # 리스트인 경우 처리
                if isinstance(data, list):
                    if len(data) > 0:
                        # 첫 번째 항목이 딕셔너리인 경우
                        if isinstance(data[0], dict):
                            logger.warning(f"LLM 응답이 리스트입니다. 첫 번째 항목을 사용합니다.")
                            data = data[0]
                        else:
                            logger.warning(f"LLM 응답이 딕셔너리가 아닌 리스트입니다: {type(data[0])}")
                            return self._fallback_info_extraction(summary)
                    else:
                        logger.warning("LLM 응답이 빈 리스트입니다.")
                        return self._fallback_info_extraction(summary)
                
                # 딕셔너리가 아닌 경우
                if not isinstance(data, dict):
                    logger.warning(f"LLM 응답이 딕셔너리가 아닙니다: {type(data)}")
                    return self._fallback_info_extraction(summary)
                    
                # 필수 키 확인 및 디폴트값 설정
                result = {
                    "title": data.get("title", "제목을 찾을 수 없습니다."),
                    "author": data.get("author", "작성자를 찾을 수 없습니다."),
                    "purpose": data.get("purpose", "목적을 찾을 수 없습니다."),
                    "summary": data.get("summary", "요약을 제공할 수 없습니다."),
                    "caption": data.get("caption", "캡션을 생성할 수 없습니다.")
                }
                logger.debug("JSON 변환 성공")
                return result
                
            except json.JSONDecodeError as e:
                logger.error(f"JSON 변환 실패: {str(e)}. 응답 내용:\n{content}")
                result = self._fallback_info_extraction(summary)
                result["error"] = f"JSON 변환 실패: {str(e)}"
                result["raw_response"] = content
                return result
                
        except Exception as e:
            import traceback
            logger.error(f"문서 정보 추출 실패: {str(e)}\n{traceback.format_exc()}")
            return self._fallback_info_extraction(summary)
    
    def _fallback_info_extraction(self, text: str) -> Dict[str, str]:
        """LLM 분석 실패 시 폴백 정보 추출"""
        try:
            logger.debug("폴백 정보 추출 사용 중...")
            # 빈 줄로 분리하여 첫 부분 추출
            paragraphs = [p for p in text.split('\n\n') if p.strip()]
            first_paragraph = paragraphs[0] if paragraphs else ""
            
            # 문서 시작 부분에서 제목 찾기
            title_candidates = []
            lines = text.split('\n')[:10]  # 처음 10줄만 확인
            
            for line in lines:
                clean_line = line.strip()
                # 제목 조건: 짧고, 특정 문자로 끝나지 않음
                if clean_line and len(clean_line) < 100 and not clean_line[-1] in ['.', ',', ';', ':']:
                    title_candidates.append(clean_line)
            
            title = title_candidates[0] if title_candidates else "제목 정보 없음"
            
            # 요약은 첫 문단 사용
            summary = first_paragraph if len(first_paragraph) < 500 else first_paragraph[:500] + "..."
            
            logger.debug("폴백 정보 추출 완료")
            return {
                "title": title,
                "author": "작성자 정보 없음",
                "purpose": "문서 목적 정보 없음",
                "summary": summary,
                "caption": title
            }
        except Exception as e:
            import traceback
            logger.error(f"폴백 정보 추출 실패: {str(e)}\n{traceback.format_exc()}")
            # 최소한의 정보 반환
            return {
                "title": "제목 추출 실패",
                "author": "작성자 정보 없음",
                "purpose": "문서 목적 정보 없음",
                "summary": "요약 추출 실패",
                "caption": "정보 추출 실패"
            }
    
    def _extract_metadata(self, file_path: str) -> Dict[str, Any]:
        """파일 메타데이터 추출"""
        p = Path(file_path)
        try:
            st = p.stat()
            return {
                "file_path": file_path,
                "file_name": p.name,
                "file_size": st.st_size,
                "created_time": str(datetime.fromtimestamp(st.st_ctime)),
                "modified_time": str(datetime.fromtimestamp(st.st_mtime)),
                "file_type": self._detect_file_type(p)
            }
        except Exception as e:
            logger.error(f"메타데이터 추출 오류: {str(e)}")
            return {
                "file_path": file_path,
                "file_name": p.name,
                "file_type": self._detect_file_type(p),
                "file_size": None,
                "created_time": None,
                "modified_time": None
            }
    
    def _detect_file_type(self, p: Path) -> str:
        """파일 유형 감지"""
        ext = p.suffix.lower()
        if ext in [".pdf"]:
            return "PDF"
        elif ext in [".docx", ".doc"]:
            return "DOCX"
        elif ext in [".pptx", ".ppt"]:
            return "PPTX"
        elif ext in [".xlsx", ".xls"]:
            return "XLSX"
        elif ext in [".txt"]:
            return "TXT"
        elif ext in [".md", ".markdown"]:
            return "MD"
        elif ext in [".rtf"]:
            return "RTF"
        elif ext in [".csv"]:
            return "CSV"
        else:
            return "Unknown"
    
    def _extract_table_data(self, pdf_path: str) -> List[Dict]:
        """PDF에서 표 구조를 추출"""
        if not PDFPLUMBER_AVAILABLE:
            logger.warning("pdfplumber 라이브러리가 설치되지 않았습니다. 표 추출 기능을 사용할 수 없습니다.")
            return []
            
        try:
            tables = []
            
            with pdfplumber.open(pdf_path) as pdf:
                for page_num, page in enumerate(pdf.pages, 1):
                    # 표 추출
                    pdf_tables = page.extract_tables()
                    
                    for table_idx, table in enumerate(pdf_tables, 1):
                        if not table or len(table) < 2:  # 헤더+데이터 최소 2행
                            continue
                        
                        # 헤더와 데이터 분리
                        headers = [str(col).strip() if col else '' for col in table[0]]
                        rows = []
                        for row in table[1:]:
                            cleaned_row = [str(cell).strip() if cell else '' for cell in row]
                            if any(cleaned_row):  # 빈 행 제외
                                rows.append(cleaned_row)
                                
                        if headers and rows:
                            # 표 제목 감지
                            title = self._detect_table_title(page, page_num, table_idx)
                            
                            tables.append({
                                'title': title,
                                'headers': headers,
                                'rows': rows,
                                'page': page_num
                            })
                            
            return tables
            
        except Exception as e:
            logger.error(f"표 추출 실패: {str(e)}")
            return []
    
    def _detect_table_title(self, page, page_num: int, table_idx: int) -> str:
        """표 제목 감지"""
        try:
            # 페이지 텍스트 추출
            page_text = page.extract_text() or ""
            lines = page_text.split('\n')
            
            # 표 제목 패턴 검색
            title_patterns = [
                r'(?:Table|표)\s*(\d+)[\.\:]\s*(.*)',  # "Table 1: Title" 패턴
                r'(?:표|Table)\s*(\d+)[^\w]*(.*)',     # "표 1 제목" 패턴
            ]
            
            # 패턴 매칭 시도
            for pattern in title_patterns:
                for line in lines:
                    match = re.search(pattern, line, re.IGNORECASE)
                    if match:
                        table_num = match.group(1)
                        if int(table_num) == table_idx:
                            return line.strip()
                        # 표 제목을 찾지 못한 경우 기본값 사용
            return f"표 {table_idx} (페이지 {page_num})"
                
        except Exception as e:
            logger.debug(f"표 제목 감지 실패: {str(e)}")
            return f"표 {table_idx} (페이지 {page_num})"
    
    def _extract_images_from_pdf(self, pdf_path: str) -> List[Dict]:
        """PDF에서 이미지 추출"""
        if not PYMUPDF_AVAILABLE:
            logger.warning("PyMuPDF(fitz) 라이브러리가 설치되지 않았습니다. 이미지 추출 기능을 사용할 수 없습니다.")
            return []
            
        try:
            images = []
            
            doc = fitz.open(pdf_path)
            for page_idx in range(len(doc)):
                page = doc[page_idx]
                image_list = page.get_images(full=True)
                
                for img_idx, img in enumerate(image_list):
                    xref = img[0]  # 이미지 참조 번호
                    
                    try:
                        base_image = doc.extract_image(xref)
                        image_bytes = base_image["image"]
                        image_ext = base_image["ext"]
                        
                        # 이미지 위치 정보
                        pix = fitz.Pixmap(doc, xref)
                        rect = page.get_image_bbox(pix)
                        
                        # 이미지 주변 텍스트 분석하여 캡션 추출 시도
                        caption = self._try_find_image_caption(page, rect)
                        
                        # 이미지 정보 저장
                        image_info = {
                            "page": page_idx + 1,
                            "index": img_idx + 1,
                            "type": image_ext,
                            "caption": caption or f"이미지 {img_idx+1} (페이지 {page_idx+1})",
                            "position": {
                                "x0": rect.x0,
                                "y0": rect.y0,
                                "x1": rect.x1,
                                "y1": rect.y1
                            },
                            # 이미지 데이터 (선택적 포함)
                            "data": base64.b64encode(image_bytes).decode('utf-8')
                        }
                        
                        images.append(image_info)
                        
                    except Exception as e:
                        logger.error(f"이미지 추출 오류 (페이지 {page_idx+1}, 이미지 {img_idx+1}): {str(e)}")
            
            doc.close()
            return images
            
        except Exception as e:
            logger.error(f"이미지 추출 실패: {str(e)}")
            return []
    
    def _try_find_image_caption(self, page, rect) -> str:
        """이미지 주변 텍스트에서 캡션 추출 시도"""
        try:
            # 이미지 아래 영역에서 텍스트 추출 (캡션은 주로 이미지 아래에 위치)
            caption_rect = fitz.Rect(rect.x0, rect.y1, rect.x1, rect.y1 + 50)
            caption_text = page.get_text("text", clip=caption_rect)
            
            # 캡션 패턴 검색
            caption_patterns = [
                r'(?:그림|Fig(?:ure)?)\s*(\d+)[\.:]?\s*(.*)',  # "그림 1: 제목" 또는 "Figure 1: Title"
                r'(?:그래프|Graph|Chart)\s*(\d+)[\.:]?\s*(.*)'  # "그래프 1: 제목" 또는 "Chart 1: Title"
            ]
            
            for pattern in caption_patterns:
                match = re.search(pattern, caption_text, re.IGNORECASE)
                if match:
                    return caption_text.strip()
            
            # 특별한 패턴이 없으면 첫 줄만 반환
            if caption_text:
                return caption_text.strip().split('\n')[0]
                
            return ""
        except Exception as e:
            logger.debug(f"캡션 추출 시도 실패: {str(e)}")
            return ""
    
    def _structure_text(self, docs: List[Document]) -> Dict[str, Any]:
        """문서를 구조화된 형태로 변환"""
        pages = []
        
        for i, d in enumerate(docs, start=1):
            # 메타데이터 정보 추출
            metadata = d.metadata.copy() if hasattr(d, 'metadata') else {}
            
            # 페이지 번호 설정
            if 'page' in metadata:
                page_num = metadata['page']
            elif 'page_number' in metadata:
                page_num = metadata['page_number']
            elif 'slide' in metadata:
                page_num = metadata['slide']
            elif 'section' in metadata:
                page_num = metadata['section']
            else:
                page_num = i
            
            # 텍스트 미리보기 생성
            content = d.page_content
            preview_len = min(300, len(content))
            preview = content[:preview_len] + ("..." if len(content) > preview_len else "")
            
            # 텍스트 분리 (문단 단위)
            paragraphs = []
            for para in content.split('\n\n'):
                if para.strip():
                    paragraphs.append(para.strip())
            
            # 페이지 정보 추가
            page_info = {
                "page_index": page_num,
                "text_preview": preview,
                "paragraphs": paragraphs[:5],  # 처음 5개 문단만 포함
                "total_paragraphs": len(paragraphs)
            }
            
            # 메타데이터 통합
            for key, value in metadata.items():
                if key not in ['source', 'page', 'page_number']:
                    page_info[key] = value
            
            pages.append(page_info)
        
        # 페이지 번호 순으로 정렬
        pages.sort(key=lambda x: x['page_index'])
        
        return {"pages": pages}
    
    def save_result_as_markdown(self, result: Dict[str, Any], output_path: str) -> str:
        """
        문서 분석 결과를 마크다운 파일로 저장
        
        Args:
            result (Dict[str, Any]): 분석 결과
            output_path (str): 저장할 경로
            
        Returns:
            str: 저장된 파일 경로
        """
        try:
            # 출력 디렉토리 생성
            output_dir = Path(output_path).parent
            output_dir.mkdir(parents=True, exist_ok=True)
            
            with open(output_path, 'w', encoding='utf-8') as f:
                # 문서 제목
                file_info = result.get('file_info', {})
                file_name = file_info.get('file_name', 'unknown')
                analysis = result.get('analysis_result', {})
                doc_title = analysis.get('title', file_name)
                
                f.write(f"# {doc_title}\n\n")
                
                # 기본 정보
                f.write("## 문서 정보\n\n")
                f.write(f"- **파일명**: {file_name}\n")
                f.write(f"- **파일 크기**: {file_info.get('file_size', 'N/A')} 바이트\n")
                f.write(f"- **생성 시간**: {file_info.get('created_time', 'N/A')}\n")
                f.write(f"- **수정 시간**: {file_info.get('modified_time', 'N/A')}\n")
                f.write(f"- **문서 길이**: {result.get('doc_length', 'N/A')} 문자\n\n")
                
                # 문서 분석 결과
                f.write("## 문서 분석\n\n")
                f.write(f"- **제목**: {analysis.get('title', 'N/A')}\n")
                f.write(f"- **작성자**: {analysis.get('author', 'N/A')}\n")
                f.write(f"- **목적**: {analysis.get('purpose', 'N/A')}\n\n")
                
                # 요약
                f.write("### 요약\n\n")
                f.write(f"{analysis.get('summary', 'N/A')}\n\n")
                
                # 표
                tables = result.get('tables', [])
                if tables:
                    f.write("## 표\n\n")
                    for i, table in enumerate(tables, 1):
                        f.write(f"### {table.get('title', f'표 {i}')}\n\n")
                        
                        # 마크다운 표 형식으로 변환
                        headers = table.get('headers', [])
                        if headers:
                            f.write('| ' + ' | '.join(headers) + ' |\n')
                            f.write('|' + '---|' * len(headers) + '\n')
                            
                            # 데이터 행
                            for row in table.get('rows', []):
                                f.write('| ' + ' | '.join(row) + ' |\n')
                        
                        f.write('\n')
                
                # 이미지
                images = result.get('images', [])
                if images:
                    f.write("## 이미지\n\n")
                    for i, image in enumerate(images, 1):
                        caption = image.get('caption', f"이미지 {i}")
                        page = image.get('page', 'N/A')
                        f.write(f"### {caption}\n\n")
                        f.write(f"- **페이지**: {page}\n")
                        f.write(f"- **유형**: {image.get('type', 'N/A')}\n\n")
                        
                        # 이미지 데이터가 있으면 마크다운에 포함
                        if 'data' in image:
                            img_data = image['data']
                            img_type = image['type'].lower()
                            f.write(f"![{caption}](data:image/{img_type};base64,{img_data})\n\n")
                
                # 전체 내용
                f.write("## 전체 내용\n\n")
                f.write(result.get('final_text_or_summary', 'N/A'))
                
                # 페이지별 내용
                pages = result.get('text_content', {}).get('pages', [])
                if pages:
                    f.write("\n\n## 페이지별 내용\n\n")
                    for page in pages:
                        page_idx = page.get('page_index', 'N/A')
                        f.write(f"### 페이지 {page_idx}\n\n")
                        f.write(page.get('text_preview', 'N/A') + "\n\n")
            
            logger.info(f"마크다운 파일 저장 완료: {output_path}")
            return output_path
            
        except Exception as e:
            logger.error(f"마크다운 파일 저장 실패: {str(e)}")
            return ""

# 메인 실행 함수
def main():
    """테스트 실행 함수"""
    # 사용 가능한 모델 표시
    valid_models = DOCUMENT_PROCESSOR_MODELS.keys()
    print(f"사용 가능한 모델: {', '.join(valid_models)}")
    
    # 모델 선택
    while True:
        selected_model = input("\n사용할 모델을 선택하세요 (openai/gemini/claude/groq): ").strip().lower()
        if selected_model in valid_models:
            break
        print("올바른 모델명을 입력하세요.")
    
    # 문서 경로 입력
    file_paths_input = input("\n분석할 문서 경로(세미콜론 ';'로 구분): ").strip()
    if not file_paths_input:
        print("문서 경로가 입력되지 않았습니다. 종료.")
        return

    # 출력 디렉토리 설정
    output_dir = input("\n결과 저장 디렉토리 (기본값: ./output): ").strip() or "./output"
    
    # 파일 경로 분리
    file_paths = [p.strip() for p in file_paths_input.split(";")]
    
    # 분석기 초기화
    analyzer = DocumentAnalyzer(selected_model)
    
    # 문서 처리
    for file_path in file_paths:
        print(f"\n====== 문서 분석 시작: {file_path} ======")
        
        # 단일 문서 처리
        result = analyzer.process_single_document(file_path)
        
        # 오류 확인
        if "error" in result:
            print(f"[오류 발생] {result['error']}")
            continue
        
        # 결과 요약 출력
        print("\n===== 분석 결과 요약 =====")
        
        # 파일 정보
        file_info = result.get("file_info", {})
        print(f"파일명: {file_info.get('file_name', 'N/A')}")
        print(f"파일 크기: {file_info.get('file_size', 'N/A')} 바이트")
        print(f"문서 유형: {file_info.get('file_type', 'N/A')}")
        print(f"문서 길이: {result.get('doc_length', 'N/A')} 문자")
        
        # 표 및 이미지 정보
        tables = result.get("tables", [])
        images = result.get("images", [])
        print(f"추출된 표: {len(tables)}개")
        print(f"추출된 이미지: {len(images)}개")
        
        # 분석 결과
        analysis = result.get("analysis_result", {})
        print("\n문서 제목:", analysis.get("title", "N/A"))
        print("문서 작성자:", analysis.get("author", "N/A"))
        print("문서 목적:", analysis.get("purpose", "N/A"))
        
        # 마크다운 파일 저장
        file_name = Path(file_path).stem
        output_path = Path(output_dir) / f"{file_name}_analysis.md"
        saved_path = analyzer.save_result_as_markdown(result, str(output_path))
        
        if saved_path:
            print(f"\n분석 결과가 다음 경로에 저장되었습니다: {saved_path}")
    
    print("\n모든 문서 처리 완료!")

if __name__ == "__main__":
    main()