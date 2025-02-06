import os
import json
import re
from typing import List, Dict, Any
from pathlib import Path
from datetime import datetime
import dotenv
import time
import random
from config import DOCUMENT_PROCESSOR_MODELS, TEMPERATURES, LLM_API_KEY
from ProcessorPrompt import DOCUMENT_PROCESSOR_PROMPT

# langchain & community
from langchain_community.document_loaders import (
    PDFPlumberLoader,
    UnstructuredExcelLoader,
    Docx2txtLoader
)
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.schema import Document, HumanMessage, SystemMessage
from langchain.chains.summarize import load_summarize_chain

# 토큰 사용량 모니터링 (옵션)
from langchain_community.callbacks import get_openai_callback

# docling + python-pptx
try:
    from docling.document_converter import DocumentConverter
    DOCLING_AVAILABLE = True
except ImportError:
    DOCLING_AVAILABLE = False

try:
    from pptx import Presentation
    PYTHON_PPTX_AVAILABLE = True
except ImportError:
    PYTHON_PPTX_AVAILABLE = False

# 로깅 설정
import logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# LLM 설정
TEMPERATURE = TEMPERATURES["document"]
SHORT_DOC_THRESHOLD = 3000  # 짧은 문서 기준 (문자 수)

class DocumentAnalyzer:
    """
    (1) 문서 로딩 (PPTX는 docling 우선 → python-pptx fallback, 기타 파일은 해당 로더 사용)
    (2) 문서 길이에 따라 처리:
        - 3000자 미만: 전체 문서 내용을 사용하여 JSON 분석 수행
        - 3000자 이상: map‐reduce 요약 결과를 사용하여 JSON 분석 수행
    (3) 최종적으로 title, author, purpose, summary, caption 추출
    """
    def __init__(self, selected_model: str = "openai"):
        self.selected_model = selected_model.lower()
        self.llm = None
        self._setup_llm()
        
    def _retry_with_backoff(self, func, args=None, kwargs=None, max_retries=3):
        """지수 백오프를 적용한 재시도"""
        args = args or []
        kwargs = kwargs or {}
        
        for attempt in range(max_retries):
            try:
                return func(*args, **kwargs)
            except Exception as e:
                if "429" in str(e) and attempt < max_retries - 1:
                    wait_time = (2 ** attempt) + random.uniform(0, 1)
                    logger.warning(f"Rate limit exceeded. Waiting {wait_time:.1f}s...")
                    time.sleep(wait_time)
                    continue
                raise

    def _setup_llm(self):
        """LLM 모델 설정 (config.py 기반, API 키 예외 처리 포함)"""
        if self.selected_model not in DOCUMENT_PROCESSOR_MODELS:
            logger.error(f"지원하지 않는 모델: {self.selected_model}")
            self.llm = None
            return

        api_model, model_class = DOCUMENT_PROCESSOR_MODELS[self.selected_model]
        api_key = LLM_API_KEY.get(self.selected_model)
        if not api_key:
            logger.error(f"API 키 누락: {self.selected_model} API 키를 .env에 설정하세요.")
            self.llm = None
            return

        try:
            self.llm = model_class(api_key=api_key, model=api_model, temperature=TEMPERATURE)
            logger.info(f"{self.selected_model} 모델 초기화 성공")
        except Exception as e:
            logger.error(f"{self.selected_model} 모델 초기화 실패: {str(e)}")
            self.llm = None

    def process_documents(self, file_paths: List[str]) -> List[Dict[str, Any]]:
        results = []
        for fp in file_paths:
            try:
                result = self.process_single_document(fp)
                results.append(result)
            except Exception as e:
                results.append({"file_path": fp, "error": str(e)})
        return results

    def _extract_table_data(self, pdf_path: str) -> List[Dict]:
        """PDF에서 표 구조를 추출"""
        try:
            import pdfplumber  # 직접 임포트
            tables = []
            
            with pdfplumber.open(pdf_path) as pdf:
                for page in pdf.pages:
                    # 페이지 크기 정보
                    bbox = page.bbox  # [x0, y0, x1, y1]
                    
                    # 표 추출
                    pdf_tables = page.extract_tables()
                    
                    for table in pdf_tables:
                        if not table or len(table) < 2:  # 헤더+데이터 최소 2행
                            continue
                        
                        # 헤더와 데이터 분리
                        headers = [str(col).strip() for col in table[0] if col]
                        rows = []
                        for row in table[1:]:
                            cleaned_row = [str(cell).strip() if cell else '' for cell in row]
                            if any(cleaned_row):  # 빈 행 제외
                                rows.append(cleaned_row)
                                
                        if headers and rows:
                            # 표 영역 정보 (page.bbox 사용)
                            table_area = {
                                'bbox': [
                                    float(bbox[0]), 
                                    float(bbox[1]), 
                                    float(bbox[2]), 
                                    float(bbox[3])
                                ],
                                'page_number': page.page_number
                            }
                            
                            # 표 제목 감지
                            title = self._detect_table_title(page, table_area)
                            
                            tables.append({
                                'title': title,
                                'headers': headers,
                                'rows': rows
                            })
                            
            return tables
            
        except Exception as e:
            logger.error(f"표 추출 실패: {str(e)}")
            return []
    
    def _detect_table_title(self, page, table_area: dict) -> str:
        """표 위치 기준으로 제목 텍스트 감지"""
        try:
            # 표 위쪽 영역 좌표 계산
            try:
                # 좌표값을 숫자로 변환
                x0 = float(table_area['bbox'][0])
                x1 = float(table_area['bbox'][2])
                y0 = float(table_area['bbox'][1])
                y1 = y0 - 50  # 50pt 위
                
                if y1 < 0:  # 페이지 상단을 넘어가지 않도록
                    y1 = 0
                    
                search_area = {
                    'x0': x0,
                    'x1': x1,
                    'y0': y1,
                    'y1': y0
                }
                
            except (ValueError, TypeError, KeyError) as e:
                logger.debug(f"좌표 변환 실패, 기본값 사용: {str(e)}")
                # 기본 검색 영역 사용
                search_area = {
                    'x0': 0,
                    'x1': page.width,
                    'y0': 0,
                    'y1': page.height
                }
            
            # 영역 내 텍스트 추출
            text = page.crop(search_area).extract_text()
            if not text:
                return ""
                
            # 가능한 제목 라인
            candidates = text.strip().split('\n')
            
            # 제목 후보 선택
            for line in reversed(candidates):
                if self._is_valid_title(line):
                    return line.strip()
                    
            return ""
            
        except Exception as e:
            logger.error(f"표 제목 감지 실패: {str(e)}")
            return ""

    def process_single_document(self, file_path: str) -> Dict[str, Any]:
        """문서 처리 및 분석"""
        CHUNK_SIZE = 2000
        CHUNK_OVERLAP = 200
        
        try:
            # 메타데이터 추출
            metadata = self._extract_metadata(file_path)
            
            # 문서 로딩
            docs = self._load_document(file_path)
            if not docs:
                return {
                    "file_info": metadata,
                    "error": "문서에서 텍스트를 추출하지 못했습니다."
                }

            # 전체 텍스트 생성
            full_text = "\n".join(d.page_content for d in docs)
            doc_length = len(full_text)
            
            # 표 추출
            tables = self._extract_table_data(file_path)
            
            # 텍스트 분할
            text_splitter = RecursiveCharacterTextSplitter(
                chunk_size=CHUNK_SIZE,
                chunk_overlap=CHUNK_OVERLAP
            )
            texts = text_splitter.split_text(full_text)
            
            # 청크별 요약 (람다 제거)
            summaries = []
            for chunk in texts:
                try:
                    summary = self._retry_with_backoff(self._generate_chunk_summary, args=[chunk])
                    if summary:
                        summaries.append(summary)
                except Exception as e:
                    logger.error(f"청크 요약 실패: {str(e)}")
            
            # 요약 통합
            if summaries:
                final_summary = self._combine_summaries(summaries)
            else:
                final_summary = full_text[:CHUNK_SIZE]
                
            # JSON 분석 (람다 제거)
            info_json = self._retry_with_backoff(
                self._extract_info_json, 
                args=[final_summary, len(texts) == 1]
            )
            
            # 페이지 구조화
            structured_text = self._structure_text(docs)

            return {
                "file_info": metadata,
                "doc_length": doc_length,
                "final_text_or_summary": final_summary,
                "tables": tables,
                "text_content": structured_text,
                "analysis_result": info_json
            }
            
        except Exception as e:
            logger.error(f"문서 처리 중 오류 발생: {str(e)}")
            return {
                "file_info": metadata,
                "error": str(e)
            }
            
    def _generate_chunk_summary(self, text: str) -> str:
        """텍스트 청크 요약 생성
        
        Args:
            text (str): 요약할 텍스트 청크
            
        Returns:
            str: 생성된 요약 또는 빈 문자열
        """
        try:
            # 프롬프트 생성
            prompt = SystemMessage(content=DOCUMENT_PROCESSOR_PROMPT["system_summary"])
            user_prompt = HumanMessage(content=DOCUMENT_PROCESSOR_PROMPT["user_summary"].format(
                full_text=text
            ))
            
            # LLM 호출
            response = self.llm.invoke([prompt, user_prompt])
            
            return response.content
            
        except Exception as e:
            logger.error(f"청크 요약 생성 실패: {str(e)}")
            return ""

    def _combine_summaries(self, summaries: List[str]) -> str:
        """청크 요약들을 하나로 통합
        
        Args:
            summaries (List[str]): 청크별 요약 리스트
            
        Returns:
            str: 통합된 요약
        """
        # 빈 요약 제거
        valid_summaries = [s for s in summaries if s.strip()]
        if not valid_summaries:
            return ""
            
        # 단일 요약이면 그대로 반환
        if len(valid_summaries) == 1:
            return valid_summaries[0]
            
        # 여러 요약 통합
        combined = "\n\n".join(valid_summaries)
        
        try:
            # 최종 요약 생성
            prompt = SystemMessage(content=DOCUMENT_PROCESSOR_PROMPT["system_summary"])
            user_prompt = HumanMessage(content=DOCUMENT_PROCESSOR_PROMPT["user_summary"].format(
                full_text=combined
            ))
            
            response = self.llm.invoke([prompt, user_prompt])
            return response.content
            
        except Exception as e:
            logger.error(f"요약 통합 실패: {str(e)}")
            return combined  # 실패시 단순 연결

    def _extract_metadata(self, file_path: str) -> Dict[str, Any]:
        p = Path(file_path)
        try:
            st = p.stat()
            return {
                "file_path": file_path,
                "file_name": p.name,
                "file_size": st.st_size,
                "created_time": str(datetime.fromtimestamp(st.st_ctime)),
                "modified_time": str(datetime.fromtimestamp(st.st_mtime)),
                "file_type": self._detect_file_type(p)
            }
        except:
            return {
                "file_path": file_path,
                "file_name": p.name,
                "file_type": self._detect_file_type(p),
                "file_size": None,
                "created_time": None,
                "modified_time": None
            }

    def _detect_file_type(self, p: Path) -> str:
        ext = p.suffix.lower()
        if ext == ".pdf":
            return "PDF"
        elif ext == ".pptx":
            return "PPTX"
        elif ext == ".xlsx":
            return "XLSX"
        elif ext == ".docx":
            return "DOCX"
        else:
            return "Unknown"

    def _load_document(self, file_path: str) -> List[Document]:
        """문서 로딩: PPTX는 docling 우선, 실패 시 python-pptx fallback / 기타 파일은 해당 로더 사용"""
        ftype = self._detect_file_type(Path(file_path))
        try:
            if ftype == "PDF":
                loader = PDFPlumberLoader(file_path)
                return loader.load()
            elif ftype == "XLSX":
                loader = UnstructuredExcelLoader(file_path)
                return loader.load()
            elif ftype == "DOCX":
                loader = Docx2txtLoader(file_path)
                return loader.load()
            elif ftype == "PPTX":
                docs = self._try_docling_pptx(file_path)
                if docs and len(docs) > 0:
                    return docs
                text = self._fallback_python_pptx(file_path)
                if not text:
                    return []
                return [Document(page_content=text, metadata={"source": file_path})]
            else:
                raise ValueError("지원하지 않는 파일 형식")
        except Exception as e:
            print(f"[에러] 문서 로딩 실패: {str(e)}")
            return []

    def _try_docling_pptx(self, file_path: str) -> List[Document]:
        if not DOCLING_AVAILABLE:
            print("[docling] not installed, skipping docling process.")
            return []
        try:
            from docling.document_converter import DocumentConverter
            converter = DocumentConverter()
            result = converter.convert(file_path)
            md_text = result.document.export_to_markdown()
            if not md_text.strip():
                print("[docling] 변환 결과가 비어있음. fallback 예정")
                return []
            return [Document(page_content=md_text, metadata={"source": file_path})]
        except Exception as e:
            print(f"[docling] 변환 실패: {str(e)}")
            return []

    def _fallback_python_pptx(self, file_path: str) -> str:
        if not PYTHON_PPTX_AVAILABLE:
            print("[python-pptx] not installed, cannot fallback.")
            return ""
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            slides_text = []
            for slide in prs.slides:
                slide_texts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_texts.append(shape.text)
                slides_text.append("\n".join(slide_texts))
            return "\n\n".join(slides_text)
        except Exception as e:
            print(f"[python-pptx] 변환 실패: {str(e)}")
            return ""

    def _split_docs(self, docs: List[Document]) -> List[Document]:
        splitter = RecursiveCharacterTextSplitter(chunk_size=500, chunk_overlap=50)
        splitted = []
        for d in docs:
            sub_texts = splitter.split_text(d.page_content)
            for t in sub_texts:
                splitted.append(Document(page_content=t, metadata=d.metadata))
        return splitted

    def _map_reduce_summary(self, splitted_docs: List[Document]) -> str:
        if not self.llm:
            return "[LLM 초기화 실패]"
        chain = load_summarize_chain(self.llm, chain_type="map_reduce")
        try:
            if self.selected_model == "openai" and get_openai_callback:
                with get_openai_callback() as cb:
                    result = chain.invoke(splitted_docs)
            else:
                result = chain.invoke(splitted_docs)
            if isinstance(result, dict):
                if "output_text" in result:
                    result = result["output_text"]
                else:
                    result = json.dumps(result, ensure_ascii=False, indent=2)
            return result
        except Exception as e:
            return f"[map_reduce 요약 실패: {str(e)}]"

    def _extract_info_json(self, final_text_or_summary: str, is_short: bool) -> Dict[str, Any]:
        """
        LLM을 통해 최종 JSON 정보를 추출.
        is_short가 True이면 전체 문서 내용을, False이면 map-reduce 요약 결과를 활용.
        """
        if not self.llm:
            return {"error": "LLM 미초기화"}
        if not isinstance(final_text_or_summary, str):
            return {"error": f"전달 데이터가 문자열이 아님. 타입: {type(final_text_or_summary)}"}
        if final_text_or_summary.startswith("[map_reduce 요약 실패"):
            return {"error": final_text_or_summary}

        # LLM 응답에 코드 블록 마크다운이 포함된 경우 이를 제거하는 전처리
        def cleanup_response(text: str) -> str:
            pattern = re.compile(r"```(?:json)?\s*([\s\S]*?)\s*```", re.MULTILINE)
            m = pattern.search(text)
            if m:
                return m.group(1).strip()
            return text.strip()

        if is_short:
            system_prompt = SystemMessage(content=DOCUMENT_PROCESSOR_PROMPT["system_full"])
            user_prompt = HumanMessage(content=DOCUMENT_PROCESSOR_PROMPT["user_full"].format(final_text=final_text_or_summary))
        else:
            system_prompt = SystemMessage(content=DOCUMENT_PROCESSOR_PROMPT["system"])
            user_prompt = HumanMessage(content=DOCUMENT_PROCESSOR_PROMPT["user"].format(final_summary=final_text_or_summary))

        try:
            resp = self.llm.invoke([system_prompt, user_prompt])
            content = cleanup_response(resp.content)
            try:
                data = json.loads(content)
                return {
                    "title": data.get("title", "제목을 찾을 수 없습니다."),
                    "author": data.get("author", "작성자를 찾을 수 없습니다."),
                    "purpose": data.get("purpose", "목적을 찾을 수 없습니다."),
                    "summary": data.get("summary", "요약을 제공할 수 없습니다."),
                    "caption": data.get("caption", "캡션을 생성할 수 없습니다.")
                }
            except json.JSONDecodeError:
                logger.error(f"JSON 변환 실패. 응답 내용:\n{content}")
                return {"error": "JSON 변환 실패", "raw_response": content}
        except Exception as e:
            return {"error": f"문서 정보 추출 실패: {str(e)}"}

    def _structure_text(self, docs: List[Document]) -> Dict[str, Any]:
        pages = []
        for i, d in enumerate(docs, start=1):
            preview = d.page_content if len(d.page_content) <= 200 else d.page_content[:200] + "..."
            pages.append({"page_index": i, "text_preview": preview})
        return {"pages": pages}


def main():
    valid_models = DOCUMENT_PROCESSOR_MODELS.keys()
    while True:
        selected_model = input("\n사용할 모델을 선택하세요 (openai/gemini/claude/groq): ").strip().lower()
        if selected_model in valid_models:
            break
        print("올바른 모델명을 입력해주세요.")
    
    analyzer = DocumentAnalyzer(selected_model)
    file_paths_input = input("\n분석할 문서 경로(세미콜론 ';'로 구분): ").strip()
    if not file_paths_input:
        print("문서 경로가 입력되지 않았습니다. 종료.")
        return

    file_paths = [p.strip() for p in file_paths_input.split(";")]
    results = analyzer.process_documents(file_paths)
    print("\n=== 분석 결과 ===\n")
    for r in results:
        if "error" in r:
            print(f"[에러 발생] {r['file_info'].get('file_path', '알 수 없음')}: {r['error']}\n")
            continue
        finfo = r["file_info"]
        doc_len = r["doc_length"]
        final_text_or_summary = r["final_text_or_summary"]
        analysis_json = r["analysis_result"]
        text_preview = r["text_content"]
        print(f"--- 파일: {finfo.get('file_name')} (길이: {doc_len} chars) ---")
        print(f"    크기(byte): {finfo.get('file_size')}, 생성일: {finfo.get('created_time')}, 수정일: {finfo.get('modified_time')}")
        if text_preview["pages"]:
            first_page = text_preview["pages"][0]
            print("\n[본문 프리뷰]")
            print(f" page_index={first_page['page_index']}, text={first_page['text_preview']}")
        else:
            print("\n[본문 없음]")
        print("\n(1) 최종 요약/전체 내용:")
        if isinstance(final_text_or_summary, dict):
            final_text_or_summary = final_text_or_summary.get("output_text", json.dumps(final_text_or_summary, ensure_ascii=False, indent=2))
        if final_text_or_summary.startswith("[map_reduce 요약 실패"):
            print(final_text_or_summary)
        else:
            print(final_text_or_summary[:300] + ("..." if len(final_text_or_summary) > 300 else ""))
        print("\n(2) JSON 분석 결과:")
        if "error" in analysis_json:
            print("   [오류]", analysis_json["error"])
        else:
            print(json.dumps(analysis_json, ensure_ascii=False, indent=2))
        print("-" * 60)
    print("\n=== 프로그램 종료 ===")


if __name__ == "__main__":
    main()
