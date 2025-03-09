# document_processor.py
import os
import json
import re
import base64
import time
import random
from typing import List, Dict, Any, Optional, Tuple, Union
from pathlib import Path
from datetime import datetime

# langchain 관련 임포트
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.schema import Document, HumanMessage, SystemMessage
from langchain_community.callbacks import get_openai_callback

from logging_manager import LoggingManager
from error_handler import ErrorHandler
from base_processor import BaseProcessor
from config import DOCUMENT_PROCESSOR_MODELS, TEMPERATURES, LLM_API_KEY
from ProcessorPrompt import DOCUMENT_PROCESSOR_PROMPT

class DocumentProcessor(BaseProcessor):
    """문서 파일 처리를 담당하는 클래스"""
    
    def __init__(self, model_name: str = "claude", auto_optimize: bool = True):
        """
        DocumentProcessor 초기화
        
        Args:
            model_name (str): 사용할 LLM 모델명
            auto_optimize (bool): 자동 최적화 사용 여부
        """
        super().__init__(model_name, auto_optimize)
        
        # 문서 분석기 초기화
        self.document_analyzer = DocumentAnalyzer(model_name, auto_optimize)
    
    def _setup_model(self) -> Any:
        """
        모델 설정 (BaseProcessor 추상 메서드 구현)
        
        Returns:
            Any: 초기화된 모델 인스턴스
        """
        if self.model_name not in DOCUMENT_PROCESSOR_MODELS:
            self.logger.error(f"지원하지 않는 모델: {self.model_name}")
            return None

        api_model, model_class = DOCUMENT_PROCESSOR_MODELS[self.model_name]
        api_key = LLM_API_KEY.get(self.model_name)

        if not api_key:
            self.logger.error(f"API 키 누락: {self.model_name} API 키를 .env에 설정하세요.")
            return None

        try:
            return model_class(api_key=api_key, model=api_model, temperature=TEMPERATURES["document"])
        except Exception as e:
            raise Exception(f"{self.model_name} 모델 초기화 실패: {str(e)}")
    
    def _process_file_internal(self, file_path: str) -> Dict[str, Any]:
        """
        문서 파일 처리 내부 로직 (BaseProcessor 추상 메서드 구현)
        
        Args:
            file_path (str): 처리할 파일 경로
            
        Returns:
            Dict[str, Any]: 처리 결과
        """
        return self.document_analyzer.process_single_document(file_path)


class DocumentAnalyzer:
    """
    고급 문서 분석기
    - 문서 특성에 기반한 최적 파서 자동 선택
    - 복합 추출 기능 (텍스트, 표, 이미지)
    - 문서 구조 보존
    - 문서 길이에 따른 적응형 처리
    """
    
    def __init__(self, selected_model: str = "claude", auto_optimize: bool = True):
        """
        DocumentAnalyzer 초기화
        
        Args:
            selected_model (str): 사용할 LLM 모델명
            auto_optimize (bool): 자동 최적화 사용 여부
        """
        self.selected_model = selected_model.lower()
        self.auto_optimize = auto_optimize
        
        # 로거 및 에러 핸들러 설정
        self.logger = LoggingManager.get_instance().get_logger("document_analyzer")
        self.error_handler = ErrorHandler.get_instance()
        
        # API 호출 관리를 위한 카운터 초기화
        self.api_call_count = 0
        self.last_call_time = 0
        
        # 캐시 저장소 초기화
        self.summary_cache = {}
        
        # 문서 길이 임계값 설정
        self.short_doc_threshold = 3000    # 3000자 미만은 짧은 문서로 간주
        self.medium_doc_threshold = 10000  # 3000-10000자는 중간 문서로 간주
        # 10000자 이상은 긴 문서로 간주
        
        # 청크 크기 설정 (문서 길이에 따라 다르게 적용)
        self.chunk_sizes = {
            "short": 2000,   # 짧은 문서
            "medium": 1500,  # 중간 문서
            "long": 1000     # 긴 문서
        }
        
        # 청크 겹침 설정
        self.chunk_overlaps = {
            "short": 200,   # 짧은 문서
            "medium": 150,  # 중간 문서
            "long": 100     # 긴 문서
        }
        
        # 파서 매니저 초기화
        self.parser_manager = DocumentParserManager()
        
        # LLM 모델 설정
        self._setup_llm()
    
    def _setup_llm(self):
        """LLM 모델 설정"""
        if self.selected_model not in DOCUMENT_PROCESSOR_MODELS:
            self.logger.error(f"지원하지 않는 모델: {self.selected_model}")
            self.llm = None
            return

        api_model, model_class = DOCUMENT_PROCESSOR_MODELS[self.selected_model]
        api_key = LLM_API_KEY.get(self.selected_model)

        if not api_key:
            self.logger.error(f"API 키 누락: {self.selected_model} API 키를 .env에 설정하세요.")
            self.llm = None
            return

        try:
            self.llm = model_class(api_key=api_key, model=api_model, temperature=TEMPERATURES["document"])
            self.logger.info(f"{self.selected_model} 모델 초기화 성공")
        except Exception as e:
            error_detail = self.error_handler.handle_error(
                e, {"model_name": self.selected_model, "operation": "llm_initialization"}
            )
            self.logger.error(f"{self.selected_model} 모델 초기화 실패: {error_detail['error']}")
            self.llm = None
    
    def process_single_document(self, file_path: str) -> Dict[str, Any]:
        """
        단일 문서 파일 처리
        
        Args:
            file_path (str): 문서 파일 경로
            
        Returns:
            Dict[str, Any]: 처리 결과
        """
        try:
            # 파일 존재 확인
            file_path_obj = Path(file_path)
            if not file_path_obj.exists():
                raise FileNotFoundError(f"문서 파일을 찾을 수 없습니다: {file_path}")
            
            # 파일 메타데이터 추출
            metadata = self._extract_metadata(file_path)
            file_type = metadata.get('file_type', 'Unknown')
            
            # 최적 파서 선택 및 문서 로드
            parser_name, parser = self.parser_manager.get_optimal_parser(file_path, file_type, self.auto_optimize)
            self.logger.info(f"선택된 파서: {parser_name}")
            
            # 선택된 파서로 문서 로드 및 구조 추출
            docs, doc_structure = parser.parse_document(file_path)
            if not docs:
                return {
                    "file_info": metadata,
                    "error": "문서에서 텍스트를 추출하지 못했습니다."
                }
            
            # 전체 텍스트 생성
            full_text = "\n".join(d.page_content for d in docs)
            doc_length = len(full_text)
            
            # 문서 길이에 따른 처리 전략 결정
            if doc_length < self.short_doc_threshold:
                doc_size_category = "short"
            elif doc_length < self.medium_doc_threshold:
                doc_size_category = "medium"
            else:
                doc_size_category = "long"
                
            self.logger.info(f"문서 길이: {doc_length}자 ({doc_size_category} 카테고리)")
            
            # 문서 길이에 따른 청크 크기 및 겹침 설정
            chunk_size = self.chunk_sizes[doc_size_category]
            chunk_overlap = self.chunk_overlaps[doc_size_category]
            
            # 텍스트 분할
            text_splitter = RecursiveCharacterTextSplitter(
                chunk_size=chunk_size,
                chunk_overlap=chunk_overlap
            )
            texts = text_splitter.split_text(full_text)
            
            # 처리할 청크 선택 (문서 길이에 따라 다르게)
            chunks_to_process = self._select_chunks_by_doc_size(texts, doc_size_category)
            
            # 선택된 청크만 요약
            summaries = self._process_chunks(chunks_to_process)
            
            # 요약 통합
            final_summary = self._generate_final_summary(summaries, full_text)
                
            # JSON 분석 결과 생성
            info_json = self._generate_info_json(final_summary, doc_length < self.short_doc_threshold)
            
            # 표와 이미지 추출 (지원되는 경우)
            tables = parser.extract_tables(file_path) if hasattr(parser, 'extract_tables') else []
            images = parser.extract_images(file_path) if hasattr(parser, 'extract_images') else []
            
            # 페이지 구조화 (한 페이지당 미리보기 생성)
            structured_text = self._structure_text(docs)
            
            # 최종 결과 구성
            result = {
                "file_info": metadata,
                "doc_length": doc_length,
                "doc_structure": doc_structure,
                "final_text_or_summary": final_summary,
                "tables": tables,
                "images": images,
                "text_content": structured_text,
                "analysis_result": info_json
            }
            
            return result
            
        except Exception as e:
            error_detail = self.error_handler.handle_error(
                e, {"file_path": file_path, "operation": "document_processing"}
            )
            self.logger.error(f"문서 처리 중 오류 발생: {error_detail['error']}")
            
            # 최소한의 결과 반환
            metadata = self._extract_metadata(file_path)
            return {
                "file_info": metadata,
                "error": str(e)
            }
    
    def _extract_metadata(self, file_path: str) -> Dict[str, Any]:
        """파일 메타데이터 추출"""
        try:
            path = Path(file_path)
            stats = path.stat()
            
            return {
                "file_path": str(file_path),
                "file_name": path.name,
                "file_size": stats.st_size,
                "created_time": str(datetime.fromtimestamp(stats.st_ctime)),
                "modified_time": str(datetime.fromtimestamp(stats.st_mtime)),
                "file_type": self._detect_file_type(path)
            }
        except Exception as e:
            error_detail = self.error_handler.handle_error(
                e, {"file_path": file_path, "operation": "metadata_extraction"}
            )
            self.logger.error(f"메타데이터 추출 오류: {error_detail['error']}")
            return {
                "file_path": str(file_path),
                "file_name": Path(file_path).name,
                "file_type": self._detect_file_type(Path(file_path)),
                "file_size": None,
                "created_time": None,
                "modified_time": None
            }
    
    def _detect_file_type(self, path: Path) -> str:
        """파일 유형 감지"""
        ext = path.suffix.lower()
        if ext in [".pdf"]:
            return "PDF"
        elif ext in [".docx", ".doc"]:
            return "DOCX"
        elif ext in [".pptx", ".ppt"]:
            return "PPTX"
        elif ext in [".xlsx", ".xls"]:
            return "XLSX"
        elif ext in [".txt"]:
            return "TXT"
        elif ext in [".md", ".markdown"]:
            return "MD"
        elif ext in [".rtf"]:
            return "RTF"
        elif ext in [".csv"]:
            return "CSV"
        else:
            return "Unknown"
    
    def _select_chunks_by_doc_size(self, texts: List[str], doc_size_category: str) -> List[str]:
        """
        문서 크기별 처리할 청크 선택
        
        Args:
            texts (List[str]): 전체 청크 목록
            doc_size_category (str): 문서 크기 카테고리
            
        Returns:
            List[str]: 선택된 청크 목록
        """
        if not texts:
            return []
            
        if doc_size_category == "short" or len(texts) <= 3:
            # 짧은 문서는 모든 청크 처리
            return texts
        elif doc_size_category == "medium":
            # 중간 문서는 주요 부분만 처리
            if len(texts) <= 5:
                return texts
            else:
                # 첫 부분, 중간 부분, 마지막 부분 선택
                return [
                    texts[0],                   # 첫 번째 청크
                    texts[len(texts) // 4],     # 1/4 지점
                    texts[len(texts) // 2],     # 중간 지점
                    texts[3 * len(texts) // 4], # 3/4 지점
                    texts[-1]                   # 마지막 청크
                ]
        else:
            # 긴 문서는 더 제한적으로 처리
            if len(texts) <= 3:
                return texts
            else:
                # 첫 부분, 중간 부분, 마지막 부분만 처리
                return [
                    texts[0],                   # 첫 번째 청크
                    texts[len(texts) // 2],     # 중간 지점
                    texts[-1]                   # 마지막 청크
                ]
    
    def _process_chunks(self, chunks: List[str]) -> List[str]:
        """
        선택된 청크 처리 (요약 생성)
        
        Args:
            chunks (List[str]): 처리할 청크 목록
            
        Returns:
            List[str]: 청크별 요약 목록
        """
        if not chunks:
            self.logger.debug("처리할 청크가 없습니다.")
            return []
        
        summaries = []
        
        for i, chunk in enumerate(chunks):
            try:
                self.logger.debug(f"청크 {i+1}/{len(chunks)} 처리 중...")
                chunk_hash = hash(chunk)
                
                # 캐시에서 요약 확인
                if chunk_hash in self.summary_cache:
                    summary = self.summary_cache[chunk_hash]
                    self.logger.debug("캐시에서 요약 로드됨")
                else:
                    # 새로 요약 생성 (재시도 로직 적용)
                    summary = self._generate_chunk_summary_with_retry(chunk)
                    if summary:
                        # 요약 캐싱
                        self.summary_cache[chunk_hash] = summary
                
                summaries.append(summary)
                
            except Exception as e:
                error_detail = self.error_handler.handle_error(
                    e, {"chunk_index": i, "operation": "chunk_processing"}
                )
                self.logger.error(f"청크 {i+1} 처리 중 오류 발생: {error_detail['error']}")
                # 오류 발생 시 빈 요약 추가
                summaries.append("요약 생성 중 오류 발생")
        
        return summaries
    
    def _manage_rate_limit(self):
        """API 호출 속도 관리"""
        current_time = time.time()
        elapsed = current_time - self.last_call_time
        
        # 최소 0.5초 대기
        if elapsed < 0.5:
            time.sleep(0.5 - elapsed)
        
        # 호출 추적
        self.api_call_count += 1
        self.last_call_time = time.time()
    
    def _generate_chunk_summary_with_retry(self, text: str, max_retries: int = 3) -> str:
        """재시도 로직이 적용된 청크 요약 생성"""
        retry_count = 0
        base_delay = 1.0
        
        while retry_count < max_retries:
            try:
                # API 호출 속도 관리
                self._manage_rate_limit()
                
                return self._generate_chunk_summary(text)
                
            except Exception as e:
                retry_count += 1
                
                # 최대 재시도 초과 시 폴백 요약 반환
                if retry_count >= max_retries:
                    self.logger.warning(f"최대 재시도 횟수 초과: {str(e)}")
                    return self._fallback_summary(text)
                
                # 지수 백오프 + 지터 적용
                wait_time = (2 ** retry_count) * base_delay + random.uniform(0, 0.1 * base_delay)
                
                # 오류 유형에 따라 대기 시간 조정
                if "429" in str(e) or "quota" in str(e).lower() or "rate limit" in str(e).lower():
                    wait_time *= 2  # 속도 제한 오류는 더 오래 대기
                    self.logger.warning(f"API 할당량 초과. {wait_time:.1f}초 대기 후 재시도 ({retry_count}/{max_retries})")
                else:
                    self.logger.warning(f"오류 발생. {wait_time:.1f}초 대기 후 재시도 ({retry_count}/{max_retries})")
                
                time.sleep(wait_time)
        
        # 모든 재시도 실패 시 폴백 요약 반환
        return self._fallback_summary(text)
    
    def _generate_chunk_summary(self, text: str) -> str:
        """청크 요약 생성"""
        if not self.llm:
            self.logger.warning("LLM이 초기화되지 않았습니다. 요약을 건너뜁니다.")
            return self._fallback_summary(text)
        
        try:
            # 요약 생성 프롬프트
            prompt = SystemMessage(content=DOCUMENT_PROCESSOR_PROMPT["system_summary"])
            user_prompt = HumanMessage(content=DOCUMENT_PROCESSOR_PROMPT["user_summary"].format(
                full_text=text
            ))
            
            # LLM 호출
            if self.selected_model in ['openai', 'gemini']:
                with get_openai_callback() as cb:
                    response = self.llm.invoke([prompt, user_prompt])
                    self.logger.debug(f"토큰 사용량: {cb.total_tokens}")
            else:
                response = self.llm.invoke([prompt, user_prompt])
                
            return response.content
            
        except Exception as e:
            error_detail = self.error_handler.handle_error(
                e, {"operation": "chunk_summary"}
            )
            self.logger.error(f"청크 요약 생성 실패: {error_detail['error']}")
            return self._fallback_summary(text)
    
    def _fallback_summary(self, text: str) -> str:
        """LLM 요약 실패 시 폴백 요약 생성"""
        # 간단한 규칙 기반 요약
        sentences = re.split(r'(?<=[.!?])\s+', text)
        
        if len(sentences) <= 5:
            return text
        
        # 첫 3문장과 마지막 2문장 추출
        summary_sentences = sentences[:3] + sentences[-2:]
        summary = ' '.join(summary_sentences)
        
        return summary + "\n(자동 추출 요약)"
    
    def _generate_final_summary(self, summaries: List[str], full_text: str) -> str:
        """최종 요약 생성"""
        # 요약이 없으면 전체 텍스트의 앞부분 사용
        if not summaries:
            self.logger.debug("요약이 없어 전체 텍스트 앞부분을 사용합니다.")
            return full_text[:min(self.short_doc_threshold, len(full_text))]
            
        # 단일 요약이면 그대로 반환
        if len(summaries) == 1:
            self.logger.debug("단일 요약을 그대로 반환합니다.")
            return summaries[0]
            
        # 요약 결합
        try:
            combined = "\n\n".join(summaries)
            
            # LLM이 없거나 결합된 요약이 짧으면 그대로 반환
            if not self.llm or len(combined) < 1000:
                self.logger.debug("LLM이 없거나 요약이 짧아 결합된 요약을 그대로 반환합니다.")
                return combined
            
            # 최종 요약 생성
            self.logger.debug("LLM을 사용하여 최종 요약 생성 중...")
            
            prompt = SystemMessage(content=DOCUMENT_PROCESSOR_PROMPT["system_summary"])
            user_prompt = HumanMessage(content=DOCUMENT_PROCESSOR_PROMPT["user_summary"].format(
                full_text=combined
            ))
            
            if self.selected_model in ['openai', 'gemini']:
                with get_openai_callback() as cb:
                    response = self.llm.invoke([prompt, user_prompt])
                    self.logger.debug(f"토큰 사용량: {cb.total_tokens}")
            else:
                response = self.llm.invoke([prompt, user_prompt])
                
            self.logger.debug("최종 요약 생성 완료")
            return response.content
            
        except Exception as e:
            error_detail = self.error_handler.handle_error(
                e, {"operation": "final_summary"}
            )
            self.logger.error(f"요약 통합 실패: {error_detail['error']}")
            
            # 결합된 요약 반환
            if isinstance(summaries, list) and summaries:
                return "\n\n".join(summaries)  # 실패 시 요약 연결
            else:
                self.logger.warning("summaries가 리스트가 아니거나 비어 있습니다.")
                return str(summaries) if summaries else "요약을 생성할 수 없습니다."
    
    def _generate_info_json(self, summary: str, is_short: bool) -> Dict[str, str]:
        """문서 정보 JSON 생성"""
        if not self.llm:
            self.logger.debug("LLM이 초기화되지 않아 기본 정보 추출을 사용합니다.")
            return self._fallback_info_extraction(summary)
        
        try:
            # LLM 응답에 코드 블록 마크다운이 포함된 경우 이를 제거하는 전처리
            def cleanup_response(text: str) -> str:
                pattern = re.compile(r"```(?:json)?\s*([\s\S]*?)\s*```", re.MULTILINE)
                m = pattern.search(text)
                if m:
                    return m.group(1).strip()
                return text.strip()
            
            # 프롬프트 생성
            self.logger.debug(f"{'전체 문서' if is_short else '요약된 문서'}에 대한 정보 추출 중...")
            
            if is_short:
                system_prompt = SystemMessage(content=DOCUMENT_PROCESSOR_PROMPT["system_full"])
                user_prompt = HumanMessage(content=DOCUMENT_PROCESSOR_PROMPT["user_full"].format(final_text=summary))
            else:
                system_prompt = SystemMessage(content=DOCUMENT_PROCESSOR_PROMPT["system"])
                user_prompt = HumanMessage(content=DOCUMENT_PROCESSOR_PROMPT["user"].format(final_summary=summary))
            
            # LLM 호출
            if self.selected_model in ['openai', 'gemini']:
                with get_openai_callback() as cb:
                    resp = self.llm.invoke([system_prompt, user_prompt])
                    self.logger.debug(f"토큰 사용량: {cb.total_tokens}")
            else:
                resp = self.llm.invoke([system_prompt, user_prompt])
                
            content = cleanup_response(resp.content)
            self.logger.debug("LLM 응답 받음, JSON 변환 시도 중...")
            
            try:
                # 응답 파싱
                data = json.loads(content)
                
                # 리스트인 경우 처리
                if isinstance(data, list) and data:
                    if isinstance(data[0], dict):
                        self.logger.warning(f"LLM 응답이 리스트입니다. 첫 번째 항목을 사용합니다.")
                        data = data[0]
                    else:
                        self.logger.warning(f"LLM 응답이 딕셔너리가 아닌 리스트입니다.")
                        return self._fallback_info_extraction(summary)
                
                # 딕셔너리가 아닌 경우
                if not isinstance(data, dict):
                    self.logger.warning(f"LLM 응답이 딕셔너리가 아닙니다: {type(data)}")
                    return self._fallback_info_extraction(summary)
                    
                # 필수 키 확인 및 디폴트값 설정
                result = {
                    "title": data.get("title", "제목을 찾을 수 없습니다."),
                    "author": data.get("author", "작성자를 찾을 수 없습니다."),
                    "purpose": data.get("purpose", "목적을 찾을 수 없습니다."),
                    "summary": data.get("summary", "요약을 제공할 수 없습니다."),
                    "caption": data.get("caption", "캡션을 생성할 수 없습니다.")
                }
                self.logger.debug("JSON 변환 성공")
                return result
                
            except json.JSONDecodeError as e:
                error_detail = self.error_handler.handle_error(
                    e, {"operation": "json_parsing", "content": content[:100]}
                )
                self.logger.error(f"JSON 변환 실패: {error_detail['error']}")
                
                result = self._fallback_info_extraction(summary)
                result["error"] = f"JSON 변환 실패: {str(e)}"
                result["raw_response"] = content[:500]  # 응답 일부만 포함
                return result
                
        except Exception as e:
            error_detail = self.error_handler.handle_error(
                e, {"operation": "info_extraction"}
            )
            self.logger.error(f"문서 정보 추출 실패: {error_detail['error']}")
            return self._fallback_info_extraction(summary)
    
    def _fallback_info_extraction(self, text: str) -> Dict[str, str]:
        """LLM 분석 실패 시 폴백 정보 추출"""
        try:
            self.logger.debug("폴백 정보 추출 사용 중...")
            # 빈 줄로 분리하여 첫 부분 추출
            paragraphs = [p for p in text.split('\n\n') if p.strip()]
            first_paragraph = paragraphs[0] if paragraphs else ""
            
            # 문서 시작 부분에서 제목 찾기
            title_candidates = []
            lines = text.split('\n')[:10]  # 처음 10줄만 확인
            
            for line in lines:
                clean_line = line.strip()
                # 제목 조건: 짧고, 특정 문자로 끝나지 않음
                if clean_line and len(clean_line) < 100 and not clean_line[-1] in ['.', ',', ';', ':']:
                    title_candidates.append(clean_line)
            
            title = title_candidates[0] if title_candidates else "제목 정보 없음"
            
            # 요약은 첫 문단 사용
            summary = first_paragraph if len(first_paragraph) < 500 else first_paragraph[:500] + "..."
            
            self.logger.debug("폴백 정보 추출 완료")
            return {
                "title": title,
                "author": "작성자 정보 없음",
                "purpose": "문서 목적 정보 없음",
                "summary": summary,
                "caption": title
            }
        except Exception as e:
            error_detail = self.error_handler.handle_error(
                e, {"operation": "fallback_extraction"}
            )
            self.logger.error(f"폴백 정보 추출 실패: {error_detail['error']}")
            
            # 최소한의 정보 반환
            return {
                "title": "제목 추출 실패",
                "author": "작성자 정보 없음",
                "purpose": "문서 목적 정보 없음",
                "summary": "요약 추출 실패",
                "caption": "정보 추출 실패"
            }
    
    def _structure_text(self, docs: List[Document]) -> Dict[str, Any]:
        """
        문서를 구조화된 형태로 변환
        
        Args:
            docs (List[Document]): 문서 객체 리스트
            
        Returns:
            Dict[str, Any]: 구조화된 텍스트 정보
        """
        pages = []
        
        for i, d in enumerate(docs, start=1):
            # 메타데이터 정보 추출
            metadata = d.metadata.copy() if hasattr(d, 'metadata') else {}
            
            # 페이지 번호 설정
            if 'page' in metadata:
                page_num = metadata['page']
            elif 'page_number' in metadata:
                page_num = metadata['page_number']
            elif 'slide' in metadata:
                page_num = metadata['slide']
            elif 'section' in metadata:
                page_num = metadata['section']
            else:
                page_num = i
            
            # 텍스트 미리보기 생성
            content = d.page_content
            preview_len = min(300, len(content))
            preview = content[:preview_len] + ("..." if len(content) > preview_len else "")
            
            # 텍스트 분리 (문단 단위)
            paragraphs = []
            for para in content.split('\n\n'):
                if para.strip():
                    paragraphs.append(para.strip())
            
            # 페이지 정보 추가
            page_info = {
                "page_index": page_num,
                "text_preview": preview,
                "paragraphs": paragraphs[:5],  # 처음 5개 문단만 포함
                "total_paragraphs": len(paragraphs)
            }
            
            # 메타데이터 통합
            for key, value in metadata.items():
                if key not in ['source', 'page', 'page_number']:
                    page_info[key] = value
            
            pages.append(page_info)
        
        # 페이지 번호 순으로 정렬
        pages.sort(key=lambda x: x['page_index'])
        
        return {"pages": pages}


class DocumentParserManager:
    """
    다양한 문서 파서를 관리하는 클래스
    """
    
    def __init__(self):
        """파서 매니저 초기화"""
        self.logger = LoggingManager.get_instance().get_logger("parser_manager")
        self.error_handler = ErrorHandler.get_instance()
        
        # 사용 가능한 파서 확인 및 초기화
        self.parsers = self._initialize_parsers()
    
    def _initialize_parsers(self) -> Dict[str, Any]:
        """
        사용 가능한 모든 파서 초기화
        
        Returns:
            Dict[str, Any]: 파서 이름과 인스턴스 매핑
        """
        parsers = {}
        
        # 텍스트 파서 (기본 파서, 항상 사용 가능)
        parsers["text"] = TextParser()
        
        # PDF 파서
        try:
            import pdfplumber
            parsers["pdfplumber"] = PDFPlumberParser()
        except ImportError:
            self.logger.info("pdfplumber 라이브러리를 찾을 수 없습니다.")
        
        try:
            import fitz  # PyMuPDF
            parsers["pymupdf"] = PyMuPDFParser()
        except ImportError:
            self.logger.info("PyMuPDF (fitz) 라이브러리를 찾을 수 없습니다.")
        
        try:
            from pdfminer.high_level import extract_text
            parsers["pdfminer"] = PDFMinerParser()
        except ImportError:
            self.logger.info("pdfminer.six 라이브러리를 찾을 수 없습니다.")
        
        try:
            import PyPDF2
            parsers["pypdf2"] = PyPDF2Parser()
        except ImportError:
            self.logger.info("PyPDF2 라이브러리를 찾을 수 없습니다.")
        
        # 문서 파서
        try:
            import docx
            parsers["docx"] = DocxParser()
        except ImportError:
            self.logger.info("python-docx 라이브러리를 찾을 수 없습니다.")
        
        # 프레젠테이션 파서
        try:
            from pptx import Presentation
            parsers["pptx"] = PptxParser()
        except ImportError:
            self.logger.info("python-pptx 라이브러리를 찾을 수 없습니다.")
        
        # 스프레드시트 파서
        try:
            import openpyxl
            parsers["openpyxl"] = OpenpyxlParser()
        except ImportError:
            self.logger.info("openpyxl 라이브러리를 찾을 수 없습니다.")
        
        self.logger.info(f"초기화된 파서: {', '.join(parsers.keys())}")
        return parsers
    
    def get_available_parsers(self, file_type: str) -> Dict[str, Any]:
        """
        특정 파일 유형에 사용 가능한 파서 목록
        
        Args:
            file_type (str): 파일 유형
            
        Returns:
            Dict[str, Any]: 사용 가능한 파서 목록
        """
        available = {}
        
        # 파일 유형별 적합한 파서 선택
        if file_type == "PDF":
            for name in ["pdfplumber", "pymupdf", "pdfminer", "pypdf2"]:
                if name in self.parsers:
                    available[name] = self.parsers[name]
        
        elif file_type == "DOCX":
            if "docx" in self.parsers:
                available["docx"] = self.parsers["docx"]
        
        elif file_type == "PPTX":
            if "pptx" in self.parsers:
                available["pptx"] = self.parsers["pptx"]
        
        elif file_type == "XLSX":
            if "openpyxl" in self.parsers:
                available["openpyxl"] = self.parsers["openpyxl"]
        
        # 모든 파일 유형에 대한 기본 파서로 텍스트 파서 추가
        available["text"] = self.parsers["text"]
        
        return available
    
    def get_optimal_parser(self, file_path: str, file_type: str, auto_optimize: bool = True) -> Tuple[str, Any]:
        """
        문서 유형에 맞는 최적의 파서 가져오기
        
        Args:
            file_path (str): 파일 경로
            file_type (str): 파일 유형
            auto_optimize (bool): 자동 최적화 사용 여부
            
        Returns:
            Tuple[str, Any]: (파서 이름, 파서 인스턴스)
        """
        try:
            # 기본 파서 매핑
            default_parsers = {
                "PDF": "pdfminer",
                "DOCX": "docx",
                "PPTX": "pptx",
                "XLSX": "openpyxl",
                "TXT": "text"
            }
            
            # 자동 최적화가 비활성화된 경우 기본 파서 반환
            if not auto_optimize:
                parser_name = default_parsers.get(file_type, "text")
                
                # 기본 파서가 사용 불가능한 경우 텍스트 파서 사용
                if parser_name not in self.parsers:
                    self.logger.warning(f"기본 파서 {parser_name}가 사용 불가능합니다. 텍스트 파서를 사용합니다.")
                    return "text", self.parsers["text"]
                    
                return parser_name, self.parsers[parser_name]
            
            # 사용 가능한 파서 확인
            available_parsers = self.get_available_parsers(file_type)
            
            # 최적 파서 선택
            if file_type == "PDF":
                # PDF 파일 특성에 기반한 파서 선택
                return self._select_optimal_pdf_parser(file_path, available_parsers)
            else:
                # 기타 파일 유형은 기본 파서 사용
                parser_name = default_parsers.get(file_type, "text")
                if parser_name in available_parsers:
                    return parser_name, available_parsers[parser_name]
            
            # 적합한 파서를 찾지 못한 경우 텍스트 파서 사용
            self.logger.warning(f"{file_type} 파일 유형에 적합한 파서를 찾을 수 없습니다. 텍스트 파서를 사용합니다.")
            return "text", self.parsers["text"]
            
        except Exception as e:
            error_detail = self.error_handler.handle_error(
                e, {"file_path": file_path, "file_type": file_type}
            )
            self.logger.error(f"파서 선택 중 오류 발생: {error_detail['error']}")
            
            # 오류 발생 시 텍스트 파서 반환
            return "text", self.parsers["text"]
    
    def _select_optimal_pdf_parser(self, file_path: str, available_parsers: Dict[str, Any]) -> Tuple[str, Any]:
        """
        PDF 파일 특성에 기반한 최적 파서 선택
        
        Args:
            file_path (str): PDF 파일 경로
            available_parsers (Dict[str, Any]): 사용 가능한 파서 목록
            
        Returns:
            Tuple[str, Any]: (파서 이름, 파서 인스턴스)
        """
        # 사용 가능한 파서가 없으면 텍스트 파서 사용
        if not available_parsers or list(available_parsers.keys()) == ["text"]:
            return "text", self.parsers["text"]
        
        # 파일 특성 분석
        is_scanned, has_tables, has_images = self._analyze_pdf_features(file_path, available_parsers)
        
        # 특성에 따른 최적 파서 선택
        if is_scanned:
            # OCR이 필요하지만 현재 구현에는 없음
            # 텍스트 레이어가 있을 수 있으므로 PyMuPDF 또는 pdfminer로 시도
            if "pymupdf" in available_parsers:
                return "pymupdf", available_parsers["pymupdf"]
            elif "pdfminer" in available_parsers:
                return "pdfminer", available_parsers["pdfminer"]
        
        if has_tables:
            # 표가 있는 문서는 pdfplumber가 좋음
            if "pdfplumber" in available_parsers:
                return "pdfplumber", available_parsers["pdfplumber"]
        
        if has_images:
            # 이미지가 있는 문서는 PyMuPDF가 좋음
            if "pymupdf" in available_parsers:
                return "pymupdf", available_parsers["pymupdf"]
        
        # 기본값: 사용 가능한 첫 번째 PDF 파서
        for parser_name in ["pdfminer", "pdfplumber", "pymupdf", "pypdf2"]:
            if parser_name in available_parsers:
                return parser_name, available_parsers[parser_name]
        
        # 모든 PDF 파서가 사용 불가능한 경우 텍스트 파서 사용
        return "text", self.parsers["text"]
    
    def _analyze_pdf_features(self, file_path: str, available_parsers: Dict[str, Any]) -> Tuple[bool, bool, bool]:
        """
        PDF 파일의 특성 분석
        
        Args:
            file_path (str): PDF 파일 경로
            available_parsers (Dict[str, Any]): 사용 가능한 파서 목록
            
        Returns:
            Tuple[bool, bool, bool]: (스캔 여부, 표 포함 여부, 이미지 포함 여부)
        """
        is_scanned = False
        has_tables = False
        has_images = False
        
        try:
            # PyMuPDF로 PDF 특성 확인
            if "pymupdf" in available_parsers:
                import fitz
                
                doc = fitz.open(file_path)
                
                # 처음 몇 페이지만 확인
                max_pages = min(5, len(doc))
                
                for page_idx in range(max_pages):
                    page = doc[page_idx]
                    
                    # 이미지 확인
                    image_list = page.get_images(full=True)
                    if len(image_list) > 0:
                        has_images = True
                    
                    # 텍스트 확인 (텍스트가 적으면 스캔 문서일 가능성)
                    text = page.get_text()
                    if len(text.strip()) < 100 and len(image_list) > 0:
                        is_scanned = True
                    
                    # 표 확인 (휴리스틱 방법: 표 관련 패턴 검색)
                    if "표 " in text or "Table " in text or "|" in text:
                        has_tables = True
                
                doc.close()
            
            # pdfplumber로 표 확인
            elif "pdfplumber" in available_parsers:
                import pdfplumber
                
                with pdfplumber.open(file_path) as pdf:
                    # 처음 몇 페이지만 확인
                    max_pages = min(5, len(pdf.pages))
                    
                    for i in range(max_pages):
                        tables = pdf.pages[i].extract_tables()
                        if tables and len(tables) > 0:
                            has_tables = True
                            break
            
            # PyPDF2로 특성 확인
            elif "pypdf2" in available_parsers:
                import PyPDF2
                
                with open(file_path, 'rb') as file:
                    reader = PyPDF2.PdfReader(file)
                    
                    # 처음 몇 페이지만 확인
                    max_pages = min(5, len(reader.pages))
                    
                    for page_idx in range(max_pages):
                        page = reader.pages[page_idx]
                        text = page.extract_text()
                        
                        # 텍스트가 적으면 스캔 문서일 가능성
                        if len(text.strip()) < 100:
                            is_scanned = True
                        
                        # 표 확인 (휴리스틱 방법)
                        if "표 " in text or "Table " in text or "|" in text:
                            has_tables = True
            
        except Exception as e:
            error_detail = self.error_handler.handle_error(
                e, {"file_path": file_path, "operation": "pdf_analysis"}
            )
            self.logger.warning(f"PDF 특성 분석 실패: {error_detail['error']}")
        
        return is_scanned, has_tables, has_images


# 문서 파서 추상 클래스
class DocumentParser:
    """문서 파서 기본 클래스"""
    
    def __init__(self):
        """파서 초기화"""
        self.logger = LoggingManager.get_instance().get_logger(self.__class__.__name__)
        self.error_handler = ErrorHandler.get_instance()
    
    def parse_document(self, file_path: str) -> Tuple[List[Document], Dict[str, Any]]:
        """
        문서 파싱 (추상 메서드)
        
        Args:
            file_path (str): 파일 경로
            
        Returns:
            Tuple[List[Document], Dict[str, Any]]: (문서 리스트, 구조 정보)
        """
        raise NotImplementedError("각 파서 클래스에서 구현해야 합니다.")
    
    def extract_tables(self, file_path: str) -> List[Dict[str, Any]]:
        """
        표 추출 (선택적 구현)
        
        Args:
            file_path (str): 파일 경로
            
        Returns:
            List[Dict[str, Any]]: 추출된 표 목록
        """
        return []
    
    def extract_images(self, file_path: str) -> List[Dict[str, Any]]:
        """
        이미지 추출 (선택적 구현)
        
        Args:
            file_path (str): 파일 경로
            
        Returns:
            List[Dict[str, Any]]: 추출된 이미지 목록
        """
        return []


    # 텍스트 파서 (기본 파서)
class TextParser(DocumentParser):
    """일반 텍스트 파일용 파서"""
    
    def parse_document(self, file_path: str) -> Tuple[List[Document], Dict[str, Any]]:
        """일반 텍스트 파일 파싱"""
        structure_info = {"type": "TEXT", "parser": "text"}
        
        try:
            # 여러 인코딩 시도
            text = None
            encodings = ['utf-8', 'cp949', 'euc-kr', 'latin1']
            
            for encoding in encodings:
                try:
                    with open(file_path, 'r', encoding=encoding) as f:
                        text = f.read()
                    break
                except UnicodeDecodeError:
                    continue
            
            if text is None:
                raise ValueError("지원되지 않는 파일 인코딩")
            
            # 구조 정보 업데이트
            lines = text.split('\n')
            paragraphs = [p for p in text.split('\n\n') if p.strip()]
            
            structure_info["lines"] = len(lines)
            structure_info["paragraphs"] = len(paragraphs)
            
            # 문서 객체 생성
            metadata = {"source": file_path}
            docs = [Document(page_content=text, metadata=metadata)]
            
            return docs, structure_info
            
        except Exception as e:
            error_detail = self.error_handler.handle_error(
                e, {"file_path": file_path, "operation": "text_parsing"}
            )
            self.logger.error(f"텍스트 파일 파싱 실패: {error_detail['error']}")
            
            # 기본 문서 반환
            error_doc = Document(
                page_content=f"파일을 읽을 수 없습니다: {error_detail['error']}",
                metadata={"source": file_path, "error": True}
            )
            return [error_doc], {"error": error_detail['error']}


# PDF 파서 구현 클래스들
class PDFMinerParser(DocumentParser):
    """PDFMiner를 사용한 PDF 파서"""
    
    def parse_document(self, file_path: str) -> Tuple[List[Document], Dict[str, Any]]:
        """PDFMiner로 PDF 파싱"""
        from pdfminer.high_level import extract_pages
        from pdfminer.layout import LTTextContainer
        
        structure_info = {"type": "PDF", "parser": "pdfminer"}
        docs = []
        
        try:
            # 페이지별 텍스트 추출
            for i, page_layout in enumerate(extract_pages(file_path)):
                page_text = ""
                
                # 레이아웃 요소 처리
                for element in page_layout:
                    if isinstance(element, LTTextContainer):
                        page_text += element.get_text()
                
                # 페이지 문서 생성
                if page_text.strip():
                    metadata = {"source": file_path, "page": i + 1}
                    docs.append(Document(page_content=page_text.strip(), metadata=metadata))
            
            # 구조 정보 업데이트
            structure_info["total_pages"] = len(docs)
            
            return docs, structure_info
            
        except Exception as e:
            error_detail = self.error_handler.handle_error(
                e, {"file_path": file_path, "operation": "pdfminer_parsing"}
            )
            self.logger.error(f"PDFMiner 파싱 실패: {error_detail['error']}")
            
            # 기본 문서 반환
            error_doc = Document(
                page_content=f"PDF 파싱 실패: {error_detail['error']}",
                metadata={"source": file_path, "error": True}
            )
            return [error_doc], {"error": error_detail['error']}


class PDFPlumberParser(DocumentParser):
    """PDFPlumber를 사용한 PDF 파서"""
    
    def parse_document(self, file_path: str) -> Tuple[List[Document], Dict[str, Any]]:
        """PDFPlumber로 PDF 파싱"""
        import pdfplumber
        
        structure_info = {"type": "PDF", "parser": "pdfplumber", "pages": []}
        docs = []
        
        try:
            with pdfplumber.open(file_path) as pdf:
                structure_info["total_pages"] = len(pdf.pages)
                
                for i, page in enumerate(pdf.pages):
                    page_text = page.extract_text() or ""
                    
                    # 페이지 정보 저장
                    page_info = {
                        "page_num": i + 1,
                        "width": page.width,
                        "height": page.height,
                        "text_length": len(page_text)
                    }
                    structure_info["pages"].append(page_info)
                    
                    # 문서 객체 생성
                    metadata = {"source": file_path, "page": i + 1, "width": page.width, "height": page.height}
                    docs.append(Document(page_content=page_text, metadata=metadata))
            
            return docs, structure_info
            
        except Exception as e:
            error_detail = self.error_handler.handle_error(
                e, {"file_path": file_path, "operation": "pdfplumber_parsing"}
            )
            self.logger.error(f"PDFPlumber 파싱 실패: {error_detail['error']}")
            
            # 기본 문서 반환
            error_doc = Document(
                page_content=f"PDF 파싱 실패: {error_detail['error']}",
                metadata={"source": file_path, "error": True}
            )
            return [error_doc], {"error": error_detail['error']}
    
    def extract_tables(self, file_path: str) -> List[Dict[str, Any]]:
        """PDFPlumber로 표 추출"""
        import pdfplumber
        
        tables = []
        
        try:
            with pdfplumber.open(file_path) as pdf:
                for page_num, page in enumerate(pdf.pages, 1):
                    # 표 추출
                    pdf_tables = page.extract_tables()
                    
                    for table_idx, table in enumerate(pdf_tables, 1):
                        if not table or len(table) < 2:  # 헤더+데이터 최소 2행
                            continue
                        
                        # 헤더와 데이터 분리
                        headers = [str(col).strip() if col else '' for col in table[0]]
                        rows = []
                        for row in table[1:]:
                            cleaned_row = [str(cell).strip() if cell else '' for cell in row]
                            if any(cleaned_row):  # 빈 행 제외
                                rows.append(cleaned_row)
                                
                        if headers and rows:
                            # 표 제목 감지
                            title = self._detect_table_title(page, page_num, table_idx)
                            
                            tables.append({
                                'title': title,
                                'headers': headers,
                                'rows': rows,
                                'page': page_num
                            })
                            
            return tables
            
        except Exception as e:
            error_detail = self.error_handler.handle_error(
                e, {"file_path": file_path, "operation": "table_extraction"}
            )
            self.logger.error(f"표 추출 실패: {error_detail['error']}")
            return []
    
    def _detect_table_title(self, page, page_num: int, table_idx: int) -> str:
        """표 제목 감지"""
        try:
            # 페이지 텍스트 추출
            page_text = page.extract_text() or ""
            lines = page_text.split('\n')
            
            # 표 제목 패턴 검색
            title_patterns = [
                r'(?:Table|표)\s*(\d+)[\.\:]\s*(.*)',  # "Table 1: Title" 패턴
                r'(?:표|Table)\s*(\d+)[^\w]*(.*)',     # "표 1 제목" 패턴
            ]
            
            # 패턴 매칭 시도
            for pattern in title_patterns:
                for line in lines:
                    match = re.search(pattern, line, re.IGNORECASE)
                    if match:
                        table_num = match.group(1)
                        if int(table_num) == table_idx:
                            return line.strip()
            
            # 표 제목을 찾지 못한 경우 기본값 사용
            return f"표 {table_idx} (페이지 {page_num})"
                
        except Exception as e:
            self.logger.debug(f"표 제목 감지 실패: {str(e)}")
            return f"표 {table_idx} (페이지 {page_num})"


class PyMuPDFParser(DocumentParser):
    """PyMuPDF를 사용한 PDF 파서"""
    
    def parse_document(self, file_path: str) -> Tuple[List[Document], Dict[str, Any]]:
        """PyMuPDF로 PDF 파싱"""
        import fitz  # PyMuPDF
        
        structure_info = {"type": "PDF", "parser": "pymupdf", "pages": []}
        docs = []
        
        try:
            doc = fitz.open(file_path)
            structure_info["total_pages"] = len(doc)
            
            for i in range(len(doc)):
                page = doc[i]
                page_text = page.get_text()
                
                # 페이지 정보
                page_info = {
                    "page_num": i + 1,
                    "width": page.rect.width,
                    "height": page.rect.height,
                    "text_length": len(page_text)
                }
                structure_info["pages"].append(page_info)
                
                # 문서 객체 생성
                metadata = {"source": file_path, "page": i + 1, "width": page.rect.width, "height": page.rect.height}
                docs.append(Document(page_content=page_text, metadata=metadata))
            
            doc.close()
            return docs, structure_info
            
        except Exception as e:
            error_detail = self.error_handler.handle_error(
                e, {"file_path": file_path, "operation": "pymupdf_parsing"}
            )
            self.logger.error(f"PyMuPDF 파싱 실패: {error_detail['error']}")
            
            # 기본 문서 반환
            error_doc = Document(
                page_content=f"PDF 파싱 실패: {error_detail['error']}",
                metadata={"source": file_path, "error": True}
            )
            return [error_doc], {"error": error_detail['error']}
    
    def extract_images(self, file_path: str) -> List[Dict[str, Any]]:
        """PyMuPDF로 이미지 추출"""
        import fitz  # PyMuPDF
        
        images = []
        
        try:
            doc = fitz.open(file_path)
            
            for page_idx in range(len(doc)):
                page = doc[page_idx]
                image_list = page.get_images(full=True)
                
                for img_idx, img in enumerate(image_list):
                    xref = img[0]  # 이미지 참조 번호
                    
                    try:
                        base_image = doc.extract_image(xref)
                        image_bytes = base_image["image"]
                        image_ext = base_image["ext"]
                        
                        # 이미지 위치 정보
                        pix = fitz.Pixmap(doc, xref)
                        rect = page.get_image_bbox(pix)
                        
                        # 이미지 주변 텍스트 분석하여 캡션 추출 시도
                        caption = self._try_find_image_caption(page, rect)
                        
                        # 이미지 정보 저장
                        image_info = {
                            "page": page_idx + 1,
                            "index": img_idx + 1,
                            "type": image_ext,
                            "caption": caption or f"이미지 {img_idx+1} (페이지 {page_idx+1})",
                            "position": {
                                "x0": rect.x0,
                                "y0": rect.y0,
                                "x1": rect.x1,
                                "y1": rect.y1
                            },
                            # 이미지 데이터 (선택적 포함)
                            "data": base64.b64encode(image_bytes).decode('utf-8')
                        }
                        
                        images.append(image_info)
                        
                    except Exception as e:
                        self.logger.error(f"이미지 추출 오류 (페이지 {page_idx+1}, 이미지 {img_idx+1}): {str(e)}")
            
            doc.close()
            return images
            
        except Exception as e:
            error_detail = self.error_handler.handle_error(
                e, {"file_path": file_path, "operation": "image_extraction"}
            )
            self.logger.error(f"이미지 추출 실패: {error_detail['error']}")
            return []
    
    def _try_find_image_caption(self, page, rect) -> str:
        """이미지 주변 텍스트에서 캡션 추출 시도"""
        import fitz  # PyMuPDF
        try:
            # 이미지 아래 영역에서 텍스트 추출 (캡션은 주로 이미지 아래에 위치)
            caption_rect = fitz.Rect(rect.x0, rect.y1, rect.x1, rect.y1 + 50)
            caption_text = page.get_text("text", clip=caption_rect)
            
            # 캡션 패턴 검색
            caption_patterns = [
                r'(?:그림|Fig(?:ure)?)\s*(\d+)[\.:]?\s*(.*)',  # "그림 1: 제목" 또는 "Figure 1: Title"
                r'(?:그래프|Graph|Chart)\s*(\d+)[\.:]?\s*(.*)'  # "그래프 1: 제목" 또는 "Chart 1: Title"
            ]
            
            for pattern in caption_patterns:
                match = re.search(pattern, caption_text, re.IGNORECASE)
                if match:
                    return caption_text.strip()
            
            # 특별한 패턴이 없으면 첫 줄만 반환
            if caption_text:
                return caption_text.strip().split('\n')[0]
                
            return ""
        except Exception as e:
            self.logger.debug(f"캡션 추출 시도 실패: {str(e)}")
            return ""


class PyPDF2Parser(DocumentParser):
    """PyPDF2를 사용한 PDF 파서"""
    
    def parse_document(self, file_path: str) -> Tuple[List[Document], Dict[str, Any]]:
        """PyPDF2로 PDF 파싱"""
        import PyPDF2
        
        structure_info = {"type": "PDF", "parser": "pypdf2", "pages": []}
        docs = []
        
        try:
            with open(file_path, 'rb') as f:
                reader = PyPDF2.PdfReader(f)
                structure_info["total_pages"] = len(reader.pages)
                
                for i, page in enumerate(reader.pages):
                    page_text = page.extract_text() or ""
                    
                    # 페이지 정보
                    page_info = {
                        "page_num": i + 1,
                        "text_length": len(page_text)
                    }
                    structure_info["pages"].append(page_info)
                    
                    # 문서 객체 생성
                    metadata = {"source": file_path, "page": i + 1}
                    docs.append(Document(page_content=page_text, metadata=metadata))
            
            return docs, structure_info
            
        except Exception as e:
            error_detail = self.error_handler.handle_error(
                e, {"file_path": file_path, "operation": "pypdf2_parsing"}
            )
            self.logger.error(f"PyPDF2 파싱 실패: {error_detail['error']}")
            
            # 기본 문서 반환
            error_doc = Document(
                page_content=f"PDF 파싱 실패: {error_detail['error']}",
                metadata={"source": file_path, "error": True}
            )
            return [error_doc], {"error": error_detail['error']}


    # 문서 파서 구현
class DocxParser(DocumentParser):
    """python-docx를 사용한 DOCX 파서"""
    
    def parse_document(self, file_path: str) -> Tuple[List[Document], Dict[str, Any]]:
        """python-docx로 DOCX 파싱"""
        import docx
        
        structure_info = {"type": "DOCX", "parser": "docx", "sections": []}
        docs = []
        
        try:
            # 문서 로드
            doc = docx.Document(file_path)
            
            # 기본 정보 추출
            structure_info["paragraphs"] = len(doc.paragraphs)
            structure_info["tables"] = len(doc.tables)
            
            # 제목 추출
            headings = []
            for i, para in enumerate(doc.paragraphs):
                if para.style.name.startswith('Heading'):
                    level = int(para.style.name.replace('Heading', '')) if para.style.name != 'Heading' else 1
                    headings.append({
                        "level": level,
                        "text": para.text,
                        "position": i
                    })
            structure_info["headings"] = headings
            
            # 섹션 단위로 분리 (제목 기준)
            sections = []
            current_section = {"title": "", "content": ""}
            
            for i, para in enumerate(doc.paragraphs):
                if para.style.name.startswith('Heading'):
                    # 기존 섹션 저장
                    if current_section["content"].strip():
                        sections.append(current_section)
                    
                    # 새 섹션 시작
                    current_section = {"title": para.text, "content": ""}
                else:
                    if para.text.strip():
                        current_section["content"] += para.text + "\n"
            
            # 마지막 섹션 저장
            if current_section["content"].strip():
                sections.append(current_section)
            
            structure_info["sections"] = sections
            
            # 섹션별 문서 생성
            for i, section in enumerate(sections):
                metadata = {
                    "source": file_path,
                    "section": i + 1,
                    "title": section["title"]
                }
                docs.append(Document(page_content=section["content"], metadata=metadata))
            
            # 섹션이 없으면 전체 내용을 하나의 문서로
            if not docs:
                full_text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
                metadata = {"source": file_path}
                docs.append(Document(page_content=full_text, metadata=metadata))
            
            return docs, structure_info
            
        except Exception as e:
            error_detail = self.error_handler.handle_error(
                e, {"file_path": file_path, "operation": "docx_parsing"}
            )
            self.logger.error(f"DOCX 파싱 실패: {error_detail['error']}")
            
            # 기본 문서 반환
            error_doc = Document(
                page_content=f"DOCX 파싱 실패: {error_detail['error']}",
                metadata={"source": file_path, "error": True}
            )
            return [error_doc], {"error": error_detail['error']}


    # 프레젠테이션 파서
class PptxParser(DocumentParser):
    """python-pptx를 사용한 PPTX 파서"""
    
    def parse_document(self, file_path: str) -> Tuple[List[Document], Dict[str, Any]]:
        """python-pptx로 PPTX 파싱"""
        from pptx import Presentation
        
        structure_info = {"type": "PPTX", "parser": "pptx", "slides": []}
        docs = []
        
        try:
            # PPTX 로드
            prs = Presentation(file_path)
            
            # 구조 정보
            structure_info["total_slides"] = len(prs.slides)
            
            # 슬라이드별 처리
            for i, slide in enumerate(prs.slides):
                slide_text = []
                slide_info = {
                    "slide_num": i + 1,
                    "shapes": len(slide.shapes)
                }
                
                # 슬라이드 텍스트 추출
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text)
                
                # 슬라이드 노트 확인
                notes_text = ""
                if slide.has_notes_slide and slide.notes_slide.notes_text_frame.text.strip():
                    notes_text = slide.notes_slide.notes_text_frame.text
                    slide_info["has_notes"] = True
                
                # 구조 정보 업데이트
                combined_text = "\n".join(slide_text)
                if notes_text:
                    combined_text += "\n\n[Notes]\n" + notes_text
                
                slide_info["text_length"] = len(combined_text)
                structure_info["slides"].append(slide_info)
                
                # 문서 객체 생성
                metadata = {
                    "source": file_path,
                    "slide": i + 1,
                    "has_notes": bool(notes_text)
                }
                docs.append(Document(page_content=combined_text, metadata=metadata))
            
            return docs, structure_info
            
        except Exception as e:
            error_detail = self.error_handler.handle_error(
                e, {"file_path": file_path, "operation": "pptx_parsing"}
            )
            self.logger.error(f"PPTX 파싱 실패: {error_detail['error']}")
            
            # 기본 문서 반환
            error_doc = Document(
                page_content=f"PPTX 파싱 실패: {error_detail['error']}",
                metadata={"source": file_path, "error": True}
            )
            return [error_doc], {"error": error_detail['error']}


    # 스프레드시트 파서
class OpenpyxlParser(DocumentParser):
    """openpyxl을 사용한 XLSX 파서"""
    
    def parse_document(self, file_path: str) -> Tuple[List[Document], Dict[str, Any]]:
        """openpyxl로 XLSX 파싱"""
        import openpyxl
        
        structure_info = {"type": "XLSX", "parser": "openpyxl", "sheets": []}
        docs = []
        
        try:
            # XLSX 로드
            wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
            
            # 구조 정보
            structure_info["total_sheets"] = len(wb.sheetnames)
            structure_info["sheet_names"] = wb.sheetnames
            
            # 시트별 처리
            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                
                # 시트 정보
                sheet_info = {
                    "name": sheet_name,
                    "max_row": sheet.max_row,
                    "max_column": sheet.max_column
                }
                structure_info["sheets"].append(sheet_info)
                
                # 시트 내용 추출
                sheet_content = []
                for row in sheet.iter_rows(values_only=True):
                    if any(cell is not None and str(cell).strip() for cell in row):
                        sheet_content.append("\t".join(str(cell) if cell is not None else "" for cell in row))
                
                # 문서 객체 생성
                sheet_text = "\n".join(sheet_content)
                metadata = {
                    "source": file_path,
                    "sheet": sheet_name,
                    "rows": sheet.max_row,
                    "columns": sheet.max_column
                }
                docs.append(Document(page_content=sheet_text, metadata=metadata))
            
            # 워크북 닫기
            wb.close()
            return docs, structure_info
            
        except Exception as e:
            error_detail = self.error_handler.handle_error(
                e, {"file_path": file_path, "operation": "xlsx_parsing"}
            )
            self.logger.error(f"XLSX 파싱 실패: {error_detail['error']}")
            
            # 기본 문서 반환
            error_doc = Document(
                page_content=f"XLSX 파싱 실패: {error_detail['error']}",
                metadata={"source": file_path, "error": True}
            )
            return [error_doc], {"error": error_detail['error']}